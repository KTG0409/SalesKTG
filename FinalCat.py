#!/usr/bin/env python3
# category_status_and_leads.py
# Build: Status (YoY) + Brand split + Sites/Items tables + Weekly combo charts (+12w forecast) + Sales/Vendor leads
# Usage example:
#   python category_status_and_leads.py --groundfish groundfish.csv --alignment alignment.csv --outdir out --company "55" --attr-groups 538,550 --vendors 4074,1260 --forecast linear
#!/usr/bin/env python3
# -*- coding: utf-8 -*-
import sys
import io

# Fix Windows console encoding for emojis
if sys.platform == 'win32':
    try:
        sys.stdout = io.TextIOWrapper(sys.stdout.buffer, encoding='utf-8', errors='replace')
        sys.stderr = io.TextIOWrapper(sys.stderr.buffer, encoding='utf-8', errors='replace')
    except:
        pass

# Your existing imports continue below...
import argparse, os, sys, math
from datetime import datetime
import time
import pandas as pd
import numpy as np
import smtplib
from email.message import EmailMessage
from pathlib import Path
from openpyxl.utils import get_column_letter

# Triggered by file addition
from watchdog.observers import Observer
from watchdog.events import FileSystemEventHandler
from threading import Timer
import configparser

# CHANGED: add bar charts and combo chart support
from openpyxl.chart import LineChart, BarChart, Reference  # type: ignore
from openpyxl.utils.dataframe import dataframe_to_rows  # type: ignore
from openpyxl.styles import Font  # type: ignore
from openpyxl.chart.text import RichText
from openpyxl.drawing.text import Paragraph, ParagraphProperties, CharacterProperties
from openpyxl.chart.shapes import GraphicalProperties 
from openpyxl.drawing.fill import SolidColorFillProperties, ColorChoice

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.dml.color import RGBColor

# Load .env configuration for email message
from dotenv import load_dotenv
BASE_DIR = Path(__file__).resolve().parent
load_dotenv(BASE_DIR / "Shrimp_Prawns" / ".env")
load_dotenv()  # fallback if the above path isn't found

# DEBUG EMAIL CONFIG
print(f"\n🔍 EMAIL DEBUG:")
print(f"  SMTP_USER: {os.getenv('SMTP_USER', 'NOT SET')}")
print(f"  SMTP_PASS: {'*' * len(os.getenv('SMTP_PASS', '')) if os.getenv('SMTP_PASS') else 'NOT SET'}")
print(f"  MAIL_TO: {os.getenv('MAIL_TO', 'NOT SET')}\n")

# ---------- Email config ----------
def _as_bool(val) -> bool:
    return str(val).strip().lower() in {"1","true","yes","y","on"}

SMTP_HOST = os.getenv("SMTP_HOST", "smtp.office365.com")
SMTP_PORT = int(os.getenv("SMTP_PORT", "587"))
SMTP_USER = os.getenv("SMTP_USER", "")
SMTP_PASS = os.getenv("SMTP_PASS", "")
MAIL_FROM = os.getenv("MAIL_FROM", SMTP_USER)
# allow comma or semicolon lists
MAIL_TO = [x.strip() for x in os.getenv("MAIL_TO", "").replace(",", ";").split(";") if x.strip()]
MAIL_CC = [x.strip() for x in os.getenv("MAIL_CC", "").replace(",", ";").split(";") if x.strip()]
MAIL_SUBJECT = os.getenv("MAIL_SUBJECT", "[Auto] Category Report")

# Accept many truthy values; default to enabled if unspecified
_raw_enabled = os.getenv("EMAIL_ENABLED", "").strip()
EMAIL_ENABLED = _as_bool(_raw_enabled) if _raw_enabled else True
FORCE_EMAIL   = _as_bool(os.getenv("FORCE_EMAIL", "0"))  # overrides EMAIL_ENABLED if set

def send_email_with_attachments(subject: str, body: str, attachments: list[str]) -> None:
    # quick visibility (no secrets)
    print(f"    [email] enabled={EMAIL_ENABLED} force={FORCE_EMAIL} "
          f"user={'set' if SMTP_USER else 'missing'} to={len(MAIL_TO)} cc={len(MAIL_CC)}")
    if not (EMAIL_ENABLED or FORCE_EMAIL):
        print("    ⛔ Email disabled by EMAIL_ENABLED (set to 1/true/yes or use FORCE_EMAIL=1). Skipping.")
        return
    if not (SMTP_USER and SMTP_PASS and MAIL_TO):
        print("    ⚠ Email not configured (need SMTP_USER, SMTP_PASS, MAIL_TO). Skipping.")
        return

    msg = EmailMessage()
    msg["Subject"] = subject
    msg["From"] = MAIL_FROM or SMTP_USER
    msg["To"] = ", ".join(MAIL_TO)
    if MAIL_CC:
        msg["Cc"] = ", ".join(MAIL_CC)
    msg.set_content(body)

    for apath in attachments or []:
        try:
            with open(apath, "rb") as f:
                data = f.read()
            name = os.path.basename(apath)
            ext = os.path.splitext(name)[1].lower()
            maintype, subtype = ("application", "octet-stream")
            if ext == ".xlsx":
                maintype, subtype = ("application", "vnd.openxmlformats-officedocument.spreadsheetml.sheet")
            elif ext == ".csv":
                maintype, subtype = ("text", "csv")
            elif ext == ".pdf":
                maintype, subtype = ("application", "pdf")
            msg.add_attachment(data, maintype=maintype, subtype=subtype, filename=name)
        except Exception as e:
            print(f"    ⚠ Could not attach {apath}: {e}")

    try:
        with smtplib.SMTP(SMTP_HOST, SMTP_PORT, timeout=30) as s:
            s.ehlo()
            s.starttls()
            s.login(SMTP_USER, SMTP_PASS)
            s.send_message(msg)
        print(f"    ✓ Email sent to {len(MAIL_TO)} recipient(s).")
    except Exception as e:
        print(f"    ❌ Email send failed: {e}")
        
# ================= USER CONFIG =================
BASE_PATH = r"C:\Users\kmor6669\OneDrive - Sysco Corporation\SBC Analytics - Category Analysis"

DEFAULTS = {
    "source": None,
    "alignment": r"C:\Users\kmor6669\OneDrive - Sysco Corporation\SBC Analytics - Category Analysis\Category_Config\Shrimp_Prawns\Alignment550.csv",
    "outdir": r"C:\Users\kmor6669\OneDrive - Sysco Corporation\SBC Analytics - Category Analysis\Category_Outputs\Shrimp_Prawns",
    "watch_dir": r"C:\Users\kmor6669\OneDrive - Sysco Corporation\SBC Analytics - Category Analysis\Category_Inputs\Shrimp_Prawns",
    
    # SITE FILTERING 
    "company": '',
    
    # Category filtering
    "attr_groups": None,
    "attr_groups": None,
    "vendors": None,
    
    # Forecast method
    "forecast_method": "linear",
    
    # Sales leads config...
    "leads_company": None,
    "leads_acct_types": 'TRS,LCC',
    "min_ytd_per_week": 20,
    
    # Vendor leads config...
    "vendor_leads_active": "Y",
    "vendor_leads_respect_site_filter": "N",
    
    # Win-back targets...
    "active_customer_weeks": 8,
    
    # Presentation...
    "create_powerpoint": "Y",
    "ppt_top_n_targets": 10,
    
    # Territory routes...
    "create_territory_routes": "Y",
    "territory_routes_top_n_dsms": 10,
    
    # ========== PRICE ZONE LEADS (NEW!) ==========
    "use_price_zone_leads": "N",  # ← Set to "Y" to enable
    "price_zone_leads_path": r"C:\Users\kmor6669\OneDrive - Sysco Corporation\Desktop\Pricing\leads_guide.csv",
}

# ---------------------------- Config Loader ----------------------------
def load_category_config(category_folder):
    """Load category-specific config from config.ini"""
    config_path = os.path.join(category_folder, "config.ini")
    
    if not os.path.exists(config_path):
        raise SystemExit(f"Config file not found: {config_path}")
    
    config = configparser.ConfigParser()
    config.read(config_path)
    
    base = category_folder
    
    category_config = {
        "alignment": os.path.join(base, config.get("Paths", "alignment")),
        "outdir": os.path.join(base, config.get("Paths", "outputs", fallback="Outputs")),
        "watch_dir": os.path.join(base, config.get("Paths", "inputs", fallback="Inputs")),
        "company": config.get("Filters", "company", fallback=None) or None,
        "attr_groups": config.get("Filters", "attr_groups", fallback=None) or None,
        "vendors": config.get("Filters", "vendors", fallback=None) or None,
        "leads_acct_types": config.get("Settings", "leads_acct_types", fallback="TRS,LCC"),
        "vendor_leads_active": config.get("Settings", "vendor_leads_active", fallback="Y"),
        "active_customer_weeks": config.getint("Settings", "active_customer_weeks", fallback=8),
        "create_powerpoint": config.get("Settings", "create_powerpoint", fallback="Y"),
        "forecast_method": config.get("Settings", "forecast_method", fallback="linear"),
        "min_ytd_per_week": config.getint("Settings", "min_ytd_per_week", fallback=20),
        "ppt_top_n_targets": config.getint("Settings", "ppt_top_n_targets", fallback=10),
        "create_territory_routes": config.get("Settings", "create_territory_routes", fallback="Y"),
        "territory_routes_top_n_dsms": config.getint("Settings", "territory_routes_top_n_dsms", fallback=10),
    }
    
    if config.has_section("Email"):
        if config.get("Email", "mail_to", fallback=""):
            os.environ["MAIL_TO"] = config.get("Email", "mail_to")
        if config.get("Email", "mail_cc", fallback=""):
            os.environ["MAIL_CC"] = config.get("Email", "mail_cc")
        if config.get("Email", "mail_subject", fallback=""):
            os.environ["MAIL_SUBJECT"] = config.get("Email", "mail_subject")
    
    return category_config
# ---------------------------- Helpers ----------------------------
# ---------- Progress Indicator ----------
def print_progress(message, step=None, total=None):
    """Print progress message with optional step counter."""
    timestamp = datetime.now().strftime("%H:%M:%S")
    if step and total:
        print(f"[{timestamp}] ({step}/{total}) {message}")
    else:
        print(f"[{timestamp}] {message}")
    sys.stdout.flush()  # Force immediate output

# ---------- Email config ----------

def clean_numeric_column(series):
    """Remove $, %, commas from numeric columns before conversion."""
    return pd.to_numeric(
        series.astype(str).str.replace("$", "").str.replace("%", "").str.replace(",", "").str.strip(),
        errors="coerce"
    ).fillna(0)
    
def get_fiscal_period_info(df):
    """
    Extract fiscal period information from the dataframe.
    
    Returns:
        tuple: (min_week, max_week, fiscal_year)
    """
    if df is None or df.empty:
        return 1, 13, "2025"
    
    min_week = int(df["Fiscal Week Number"].min()) if "Fiscal Week Number" in df.columns else 1
    max_week = int(df["Fiscal Week Number"].max()) if "Fiscal Week Number" in df.columns else 13
    
    # Get the fiscal year for the latest week
    if "Fiscal Year Key" in df.columns and "Fiscal Week Number" in df.columns:
        max_fy = df[df["Fiscal Week Number"] == max_week]["Fiscal Year Key"].max()
    else:
        max_fy = "2025"
    
    return min_week, max_week, max_fy
    
def _clean_str(x):
    """Clean string but preserve leading zeros."""
    if pd.isna(x): return ""
    return str(x).strip()

def _clean_str_numeric(x):
    """Clean string and strip leading zeros for numeric ID comparison only."""
    if pd.isna(x): return ""
    s = str(x).strip()
    if s.isdigit():  # Only strip zeros if it's purely numeric
        return s.lstrip("0") or "0"
    return s

def concat_address(addr1, addr2):
    a1 = str(addr1).strip() if pd.notna(addr1) else ""
    a2 = str(addr2).strip() if pd.notna(addr2) else ""
    return (a1 + " " + a2).strip()

def pct(n, d):
    if d == 0 or pd.isna(d): return np.nan
    return n / d

def safe_div(n, d):
    try:
        return float(n) / float(d) if float(d) != 0 else np.nan
    except Exception:
        return np.nan

# Helper for Filtered Company Number/Name  
def _highlight_filtered_company_row(ws, header_row_0idx, company_name):
    """Highlight the filtered company row in yellow."""
    from openpyxl.styles import PatternFill, Font
    
    header_row = header_row_0idx + 1
    
    # Find "Company Name" column
    company_col = None
    for cell in ws[header_row]:
        if cell.value == "Company Name":
            company_col = cell.col_idx
            break
    
    if not company_col:
        return
    
    # Find and highlight matching row
    for row in ws.iter_rows(min_row=header_row+1, max_row=ws.max_row):
        cell_value = str(row[company_col-1].value).strip().lower()
        if cell_value == company_name.lower().strip():
            for cell in row:
                cell.fill = PatternFill(start_color="FFFF99", end_color="FFFF99", fill_type="solid")
                cell.font = Font(bold=True)
            break
    
# Helper that starts watcher
def start_watcher(watch_dir: str, alignment_path: str, outdir_path: str,
                  quiet_seconds: int = 8, forecast_method: str = "linear"):
    os.makedirs(watch_dir, exist_ok=True)
    os.makedirs(outdir_path, exist_ok=True)

    print("🔍 Starting file watcher service…")
    print(f"👀 Watching: {watch_dir}")
    print(f"📤 Output:  {outdir_path}")
    print(f"📎 File types: {', '.join(sorted(ALLOW_EXT))}")
    print(f"⏱ Debounce: {quiet_seconds} seconds")
    print("🟢 Waiting for files (Ctrl+C to stop)…")

    observer = Observer()
    handler = DebouncedHandler(alignment_path, outdir_path, quiet_seconds, forecast_method)
    observer.schedule(handler, watch_dir, recursive=False)
    observer.start()
    try:
        while True:
            import time as _t; _t.sleep(1)
    except KeyboardInterrupt:
        print("\n🛑 Stopping watcher…")
        observer.stop()
    observer.join()
    print("✅ Watcher stopped.")

# ---------------------------- Core (your originals, kept) ----------------------------
def load_and_prepare(source_path, alignment_path, company=None, attr_groups=None, vendors=None):
    print_progress("🔄 Loading source data...", 1, 5)
    g = pd.read_csv(source_path, dtype=str)
    
    print_progress("🔄 Loading alignment data...", 2, 5)
    a = pd.read_csv(alignment_path, dtype=str)

    def _clean_headers(df):
        cols = df.columns.str.replace(r"\s+", " ", regex=True).str.strip()
        df.columns = cols
        return df

    print_progress("🔄 Cleaning and preparing data...", 3, 5)
    g = _clean_headers(g).fillna("")
    a = _clean_headers(a).fillna("")

    # numerics used downstream
    for col in ["Pounds CY", "Pounds PY", "Fiscal Week Number"]:
        if col in g.columns:
            g[col] = pd.to_numeric(g[col], errors="coerce").fillna(0)

    # Clean but PRESERVE leading zeros for alignment keys
    for col in ["Company Number", "Lot Number", "Item Number", "True Vendor Number"]:
        if col in g.columns:
            g[col] = g[col].map(_clean_str)

    # Alignment key in alignment file
    if "Alignment Key" not in a.columns:
        raise SystemExit("Alignment file must contain 'Alignment Key'.")
    a["Alignment Key"] = a["Alignment Key"].astype(str).str.strip()

    # SUPC and SUVC preserve leading zeros
    for col in ["SUPC", "SUVC"]:
        if col not in a.columns:
            raise SystemExit(f"Alignment file missing required column '{col}'.")
        a[col] = a[col].map(_clean_str)

    # Build alignment key in source
    if "Alignment Key" in g.columns and g["Alignment Key"].astype(str).str.len().gt(0).any():
        g["Alignment Key"] = g["Alignment Key"].astype(str).str.strip()
    else:
        need = [c for c in ["Company Number", "Lot Number"] if c not in g.columns]
        if need:
            raise SystemExit(f"Source missing columns to build Alignment Key: {need}")
        g["Alignment Key"] = g["Company Number"].map(_clean_str) + g["Lot Number"].map(_clean_str)

    # Address concat for leads routing
    g["Customer Street Address"] = g.apply(
        lambda r: concat_address(r.get("Customer Address", ""), r.get("Customer Address 2", "")),
        axis=1
    )

    # Prepare alignment data
    align_keep = [c for c in [
        "Alignment Key", "Lot Number", "Lot Description",
        "Supplier Name", "SUPC", "SUVC",
        "Product Category", "Next events", "OPCO Name", "Award Volume Annualized" 
    ] if c in a.columns]
    for req in ["Alignment Key", "SUPC", "SUVC", "Supplier Name"]:
        if req not in align_keep:
            raise SystemExit(f"Alignment join cannot proceed; missing '{req}'.")

    align = a[align_keep].drop_duplicates("Alignment Key")
    
    # MERGE FIRST (before any filtering)
    print(f"\n🔗 Merging alignment data...")
    print(f"   Source rows: {len(g):,}")
    print(f"   Alignment keys: {len(align):,}")
    
    g = g.merge(align, on="Alignment Key", how="left", suffixes=("", "_ALN"))

    if "SUPC" not in g.columns or g["SUPC"].isna().all():
        print(f"\n❌ NO MATCHES FOUND!")
        print(f"Sample source keys: {g['Alignment Key'].head(5).tolist()}")
        print(f"Sample alignment keys: {align['Alignment Key'].head(5).tolist()}")
        raise SystemExit("No Alignment matches found (all SUPC null). Check keys between files.")

    matches = g["SUPC"].notna().sum()
    print(f"   ✅ Matched: {matches:,} records")

    # NOW FILTER by company (after merge, with smart comparison)
    if company:
        print(f"\n🎯 Applying company filter: '{company}'")
        before_filter = len(g)
        
        # Normalize for comparison (strip leading zeros)
        filter_normalized = _clean_str_numeric(company).lower()
        
        # Create temp normalized columns for comparison
        g["_CompanyNum_Norm"] = g["Company Number"].map(_clean_str_numeric).str.lower()
        g["_CompanyName_Norm"] = g.get("Company Name", "").astype(str).str.strip().str.lower()
        
        # Filter using normalized comparison
        g = g[
            (g["_CompanyNum_Norm"] == filter_normalized) |
            (g["_CompanyName_Norm"] == company.strip().lower())
        ].copy()
        
        # Drop temp columns
        g = g.drop(columns=["_CompanyNum_Norm", "_CompanyName_Norm"])
        
        print(f"   Filtered: {before_filter:,} → {len(g):,} rows")
        
        if len(g) == 0:
            available = pd.read_csv(source_path, dtype=str, nrows=1000)["Company Number"].unique()[:10]
            print(f"\n❌ No data found for company '{company}'")
            print(f"Available company numbers (sample): {available.tolist()}")
            raise SystemExit(f"Company filter '{company}' resulted in zero rows.")
    
    # Filter by attribute groups if specified
    if attr_groups and "Attribute Group ID" in g.columns:
        keep_ag = set(_clean_str(x) for x in attr_groups.split(","))
        before = len(g)
        g = g[g["Attribute Group ID"].map(_clean_str).isin(keep_ag)]
        print(f"   Attribute group filter: {before:,} → {len(g):,} rows")

    # Verify required columns
    for req in ["SUPC", "SUVC", "Supplier Name"]:
        if req not in g.columns:
            raise SystemExit(f"Post-join missing required column '{req}'.")

    # Calculate alignment flags (preserve leading zeros in comparison)
    g["IsAlignedItem"] = (g.get("Item Number", "").map(_clean_str) == g["SUPC"].map(_clean_str)).astype(int)
    g["IsAlignedVendor"] = (g.get("True Vendor Number", "").map(_clean_str) == g["SUVC"].map(_clean_str)).astype(int)

    # ==========================================
    # EXTRACT PRICE ZONE SUFFIX (NEW!)
    # ==========================================
    if 'Price Zone ID' in g.columns:
        print("   🔧 Extracting zone suffix from Price Zone ID...")
        
        # Format: "001-1" where 001=company, 1=zone
        g['Zone_Suffix'] = g['Price Zone ID'].astype(str).str.split('-').str[-1]
        g['Zone_Suffix_Numeric'] = pd.to_numeric(g['Zone_Suffix'], errors='coerce')
        
        print(f"      ✅ Extracted zones: {g['Zone_Suffix_Numeric'].nunique()} unique zones")
    else:
        print("   ⚠️  'Price Zone ID' column not found - zone analysis unavailable")
        g['Zone_Suffix_Numeric'] = None
    
    # ==========================================
    # CREATE COMPANY COMBO KEY (NEW!)
    # Must match PZone format for merging
    # ==========================================
    if 'Company Name' in g.columns and 'Attribute Group ID' in g.columns:
        print("   🔧 Creating Company_Combo_Key...")
        
        # Get company prefix from alignment (same logic as PZone)
        # Assuming you have "Sysco Company Prefix Code" or similar
        if 'Sysco Company Prefix Code' in g.columns:
            g['Company_Prefix'] = g['Sysco Company Prefix Code'].astype(str).str.strip()
        else:
            # Fallback: use first 3 chars of Company Name
            g['Company_Prefix'] = g['Company Name'].astype(str).str[:3].str.upper()
        
        # Create combo key: CompanyName_AttributeGroup_Prefix
        # Example: "Detroit_CD Bar & Grill_538_CPA"
        g['Company_Combo_Key'] = (
            g['Company Name'].astype(str).str.strip() + '_' +
            g.get('NPD Cuisine Type', g.get('Customer Account Type Code', '')).astype(str).str.strip() + '_' +
            g['Attribute Group ID'].astype(str).str.strip() + '_' +
            g['Company_Prefix'].astype(str).str.strip()
        )
        
        print(f"      ✅ Created combo keys: {g['Company_Combo_Key'].nunique()} unique combos")
    else:
        print("   ⚠️  Cannot create Company_Combo_Key - missing required columns")
        g['Company_Combo_Key'] = None
    
    return g, vendors

def compute_windows(df):
    # current week
    current_week = int(pd.to_numeric(df.get("Fiscal Week Number", 0), errors="coerce").fillna(0).max()) if len(df) else 0

    # Ensure the account-type column exists
    if "Customer Account Type Code" not in df.columns:
        df["Customer Account Type Code"] = "UNKNOWN"

    # REQUIRED grouping keys (DEDENT THIS - same level as the if statement above)
    key = [
        "Company Number","Company Name",
        "Customer Name","Company Customer Number",
        "Customer Account Type Code",
        "Customer DSM",        # ← ADD THIS
        "Customer Territory",
        "Lot Number","Lot Description",
        "Attribute Group ID","Attribute Group Name",
        "Business Center Name",
        "Alignment Key",
        "SUPC","SUVC","Supplier Name"
    ]
    
    missing = [c for c in key if c not in df.columns]
    if missing:
        raise SystemExit(f"compute_windows(): required columns missing: {missing}")

    df = df.copy()
    for col in ["Pounds CY","Pounds PY"]:
        if col not in df.columns: df[col] = 0.0
    # Aligned flags & derived lbs
    both_flag = ((df.get("IsAlignedItem", 0).astype(int)) & (df.get("IsAlignedVendor", 0).astype(int))).astype(int)
    df["ItemAligned_Lbs"] = df["Pounds CY"] * df.get("IsAlignedItem", 0)
    df["VendorAligned_Lbs"] = df["Pounds CY"] * df.get("IsAlignedVendor", 0)
    df["ItemVendorAligned_Lbs"] = df["Pounds CY"] * both_flag

    # === STEP 1: Calculate DISTINCT CUSTOMERS at customer level (before item-level grouping) ===
    # Aggregate to customer level first (sum across all items per customer)
    customer_level = df.groupby(
        ["Company Number", "Company Name", "Company Customer Number", "Customer Account Type Code"],
        dropna=False
    ).agg(
        Customer_Pounds_CY=("Pounds CY", "sum"),
        Customer_Pounds_PY=("Pounds PY", "sum")
    ).reset_index()

    # Count ONLY customers who actually bought (pounds > 0)
    customers_cy = customer_level[customer_level["Customer_Pounds_CY"] > 0].groupby(
        ["Company Number", "Company Name", "Customer Account Type Code"], dropna=False
    ).agg(
        Distinct_Customers_CY=("Company Customer Number", "nunique")
    ).reset_index()

    customers_py = customer_level[customer_level["Customer_Pounds_PY"] > 0].groupby(
        ["Company Number", "Company Name", "Customer Account Type Code"], dropna=False
    ).agg(
        Distinct_Customers_PY=("Company Customer Number", "nunique")
    ).reset_index()

    # Merge customer counts
    customer_counts = customers_cy.merge(
        customers_py,
        on=["Company Number", "Company Name", "Customer Account Type Code"],
        how="outer"
    ).fillna(0)

    # === STEP 2: Original YTD aggregates at item level ===
    # YTD aggregates (now per-account-type) - this is the ORIGINAL aggregation
    ytd = df.groupby(key, dropna=False).agg(
        Pounds_CY=("Pounds CY","sum"),
        Pounds_PY=("Pounds PY","sum"),
        ItemAligned_Lbs=("ItemAligned_Lbs","sum"),
        VendorAligned_Lbs=("VendorAligned_Lbs","sum"),
        ItemVendorAligned_Lbs=("ItemVendorAligned_Lbs","sum"),
    ).reset_index()

    # === STEP 3: Merge customer counts into the main dataset ===
    ytd = ytd.merge(
        customer_counts,
        on=["Company Number", "Company Name", "Customer Account Type Code"],
        how="left"
    ).fillna(0)

    # === STEP 4: Calculate deltas and derived metrics ===
    ytd["Delta_YoY_Lbs"] = ytd["Pounds_CY"] - ytd["Pounds_PY"]
    ytd["YoY_Pct"] = ytd.apply(lambda r: (r["Delta_YoY_Lbs"] / r["Pounds_PY"]) if r["Pounds_PY"] else np.nan, axis=1)

    # Customer metrics
    ytd["Delta_Customers"] = ytd["Distinct_Customers_CY"] - ytd["Distinct_Customers_PY"]
    ytd["Customer_Retention_Pct"] = ytd.apply(
        lambda r: (r["Distinct_Customers_CY"] / r["Distinct_Customers_PY"]) if r["Distinct_Customers_PY"] > 0 else np.nan, 
        axis=1
    )
    ytd["Avg_Lbs_per_Customer_CY"] = ytd.apply(
        lambda r: (r["Pounds_CY"] / r["Distinct_Customers_CY"]) if r["Distinct_Customers_CY"] > 0 else 0,
        axis=1
    )
    ytd["Avg_Lbs_per_Customer_PY"] = ytd.apply(
        lambda r: (r["Pounds_PY"] / r["Distinct_Customers_PY"]) if r["Distinct_Customers_PY"] > 0 else 0,
        axis=1
    )
    
    # Week-over-week (also per-account-type now)
    lw = df[df["Fiscal Week Number"] == current_week].groupby(key)["Pounds CY"].sum().rename("W_Lbs").reset_index() if current_week else pd.DataFrame(columns=key+["W_Lbs"])
    pw = df[df["Fiscal Week Number"] == current_week - 1].groupby(key)["Pounds CY"].sum().rename("Wm1_Lbs").reset_index() if current_week else pd.DataFrame(columns=key+["Wm1_Lbs"])
    wow = lw.merge(pw, on=key, how="outer") if len(lw) or len(pw) else pd.DataFrame(columns=key+["W_Lbs","Wm1_Lbs"])
    if "W_Lbs" not in wow.columns: wow["W_Lbs"] = 0.0
    if "Wm1_Lbs" not in wow.columns: wow["Wm1_Lbs"] = 0.0
    wow["WoW_Delta_Lbs"] = wow["W_Lbs"] - wow["Wm1_Lbs"]
    wow["WoW_Pct"] = wow.apply(lambda r: (r["WoW_Delta_Lbs"] / r["Wm1_Lbs"]) if r["Wm1_Lbs"] else np.nan, axis=1)

    # 4w vs prior 4w (also per-account-type)
    last4 = df[df["Fiscal Week Number"].between(max(current_week-3, 0), current_week, inclusive="both")] if current_week else df.iloc[0:0]
    prior4 = df[df["Fiscal Week Number"].between(max(current_week-7, 0), max(current_week-4, 0), inclusive="both")] if current_week else df.iloc[0:0]
    l4 = last4.groupby(key)["Pounds CY"].sum().rename("L4_Lbs").reset_index() if len(last4) else pd.DataFrame(columns=key+["L4_Lbs"])
    p4 = prior4.groupby(key)["Pounds CY"].sum().rename("P4_Lbs").reset_index() if len(prior4) else pd.DataFrame(columns=key+["P4_Lbs"])
    m4 = l4.merge(p4, on=key, how="outer") if len(l4) or len(p4) else pd.DataFrame(columns=key+["L4_Lbs","P4_Lbs"])
    if "L4_Lbs" not in m4.columns: m4["L4_Lbs"] = 0.0
    if "P4_Lbs" not in m4.columns: m4["P4_Lbs"] = 0.0
    m4["L4_vs_P4_Delta"] = m4["L4_Lbs"] - m4["P4_Lbs"]
    m4["L4_vs_P4_Pct"] = m4.apply(lambda r: (r["L4_vs_P4_Delta"] / r["P4_Lbs"]) if r["P4_Lbs"] else np.nan, axis=1)

    status = ytd.merge(wow, on=key, how="left").merge(m4, on=key, how="left").fillna(0)
    return status, current_week


def _format_sheet_by_headers(ws, number_headers=None, percent_headers=None):
    if ws.max_row < 2:
        return
    number_headers  = set(number_headers or [])
    percent_headers = set(percent_headers or [])

    # map header -> col index
    header_to_col = {str(c.value).strip(): c.col_idx for c in ws[1] if c.value is not None}

    def _coerce_number(cell):
        v = cell.value
        if v is None or isinstance(v, (int, float)):
            return
        s = str(v).strip().replace(",", "")
        # if someone already wrote "12345%" into a number column, strip the % and do NOT /100
        if s.endswith("%"):
            s = s[:-1]
        try:
            cell.value = float(s)
        except Exception:
            pass  # leave as-is

    def _coerce_percent(cell):
        v = cell.value
        if v is None:
            return
        if isinstance(v, (int, float)):
            # assume already fraction (0.23 -> 23%)
            return
        s = str(v).strip().replace(",", "")
        try:
            if s.endswith("%"):
                cell.value = float(s[:-1]) / 100.0
            else:
                cell.value = float(s)  # treat as fraction already
        except Exception:
            pass

    # clear any old formats on targeted cols first
    for h in (number_headers | percent_headers):
        col = header_to_col.get(h)
        if not col:
            continue
        for row in ws.iter_rows(min_row=2, max_row=ws.max_row, min_col=col, max_col=col):
            row[0].number_format = "General"

    # numbers
    for h in number_headers:
        col = header_to_col.get(h)
        if not col:
            continue
        for row in ws.iter_rows(min_row=2, max_row=ws.max_row, min_col=col, max_col=col):
            _coerce_number(row[0])
            row[0].number_format = "#,##0"

    # percents
    for h in percent_headers:
        col = header_to_col.get(h)
        if not col:
            continue
        for row in ws.iter_rows(min_row=2, max_row=ws.max_row, min_col=col, max_col=col):
            _coerce_percent(row[0])
            row[0].number_format = "0.0%"


def _try_format(xw, sheet_name, number_headers=None, percent_headers=None):
    wb = xw.book
    if sheet_name not in wb.sheetnames:
        return
    _format_sheet_by_headers(wb[sheet_name], number_headers, percent_headers)


def classify_conversion(status, X=0.80, Y=0.80, Z=0.95):
    status["Frac_ItemAligned"] = status.apply(lambda r: pct(r["ItemAligned_Lbs"], r["Pounds_CY"]), axis=1)
    status["Frac_VendorAligned"] = status.apply(lambda r: pct(r["VendorAligned_Lbs"], r["Pounds_CY"]), axis=1)
    status["Frac_ItemVendorAligned"] = status.apply(lambda r: pct(r["ItemVendorAligned_Lbs"], r["Pounds_CY"]), axis=1)

    def label(row):
        cy = row["Pounds_CY"]
        f_item = row["Frac_ItemAligned"] or 0
        f_vendor = row["Frac_VendorAligned"] or 0
        f_both = row["Frac_ItemVendorAligned"] or 0
        if cy <= 0:
            return "No CY Volume"
        if f_both >= Z:
            return "Converted"
        if f_item >= Y and f_vendor < X:
            return "Needs Vendor"
        if f_vendor >= Y and f_item < X:
            return "Needs Item"
        if f_item < X and f_vendor < X:
            return "Needs Both"
        return "Partial"

    status["Conversion_Status"] = status.apply(label, axis=1)
    return status

def mark_exit_lapsing(status, min_ytd):
    status["Is_Exit"] = ((status["Pounds_CY"] <= 0) & (status["Pounds_PY"] >= min_ytd)).astype(int)
    status["Is_Lapsing"] = ((status["Pounds_CY"] <= 0.2*status["Pounds_PY"]) & (status["Pounds_PY"] >= min_ytd)).astype(int)
    return status

def build_sales_leads(status, raw_df, min_ytd):
    # Filter to meaningful volume & actionable statuses
    leads = status[
        (status["Pounds_CY"] >= min_ytd) &
        (status["Conversion_Status"].isin(["Needs Both","Needs Item","Needs Vendor","Partial"]))
    ].copy()

    key = ["Company Number","Company Name","Customer Name","Company Customer Number","Lot Number","Lot Description","Alignment Key"]
    cols_route = ["Customer DSM","Customer Territory","Customer Street Address","Customer City","Customer State","Customer Zip Code"]
    cols_purchase = ["Item Number","Item Description","True Vendor Number","True Vendor Name","Pounds CY"]
    attach_cols = key + cols_route + cols_purchase

    latest = raw_df.copy()
    latest = latest.sort_values(["Company Customer Number","Lot Number","Fiscal Week Number"], ascending=[True,True,False])
    latest = latest.drop_duplicates(subset=key, keep="first")

    slim = leads[key + ["Pounds_CY","Conversion_Status","SUPC","SUVC","Supplier Name","Frac_ItemAligned","Frac_VendorAligned","Frac_ItemVendorAligned"]]
    out = slim.merge(latest[attach_cols], on=key, how="left")

    def what_to_convert(row):
        if row["Conversion_Status"] == "Needs Both": return "Convert Item + Vendor"
        if row["Conversion_Status"] == "Needs Item": return "Convert Item"
        if row["Conversion_Status"] == "Needs Vendor": return "Convert Vendor"
        if row["Conversion_Status"] == "Partial": return "Top off to aligned"
        return ""
    out["Action"] = out.apply(what_to_convert, axis=1)

    out.rename(columns={"Supplier Name":"Aligned Supplier Name"}, inplace=True)
    out["Aligned Item (SUPC)"] = out["SUPC"]
    out["Aligned Vendor (SUVC)"] = out["SUVC"]

    ordered = [
        "Company Number","Company Name","Customer Name","Company Customer Number",
        "Customer DSM","Customer Territory","Customer Street Address","Customer City","Customer State","Customer Zip Code",
        "Lot Number","Lot Description","Alignment Key","Aligned Supplier Name","Aligned Item (SUPC)","Aligned Vendor (SUVC)",
        "Item Number","Item Description","True Vendor Number","True Vendor Name",
        "Pounds_CY","Frac_ItemVendorAligned","Frac_ItemAligned","Frac_VendorAligned",
        "Conversion_Status","Action"
    ]
    existing = [c for c in ordered if c in out.columns]
    remaining = [c for c in out.columns if c not in existing]
    out = out[existing + remaining]
    return out

def vendor_splits(leads_df, vendors_filter, outdir):
    """
    Create individual CSV files for each vendor with their TRS-only leads.
    
    NOTE: Vendor splits are ALWAYS filtered to TRS accounts only.
    
    Args:
        leads_df: Full leads dataframe
        vendors_filter: Comma-separated vendor numbers to include (None = all)
        outdir: Output directory
    """
    vendor_dir = os.path.join(outdir, "Vendor_Leads")
    os.makedirs(vendor_dir, exist_ok=True)
    
    if "Aligned Vendor (SUVC)" not in leads_df.columns:
        print("    ⚠ No 'Aligned Vendor (SUVC)' column - skipping vendor splits")
        return []
    
    # HARDCODED: Vendor splits are ALWAYS TRS only
    subset = leads_df.copy()
    if "Customer Account Type Code" in subset.columns:
        subset = subset[subset["Customer Account Type Code"] == "TRS"]
        print(f"    Vendor splits: TRS only ({len(subset)} leads)")
    else:
        print(f"    ⚠ Customer Account Type Code not found - cannot filter to TRS")
        return []
    
    # Get valid vendors
    subset = subset[subset["Aligned Vendor (SUVC)"].notna() & (subset["Aligned Vendor (SUVC)"]!="")]
    
    if subset.empty:
        print("    ℹ No TRS leads found for vendor splits")
        return []
    
    aligned_vendors = subset["Aligned Vendor (SUVC)"].dropna().unique().tolist()
    
    # Apply vendor filter if specified
    if vendors_filter:
        keep = set([v.strip() for v in vendors_filter.split(",") if v.strip()])
        aligned_vendors = [v for v in aligned_vendors if v in keep]
        print(f"    Vendor splits filtered to vendors: {keep}")
    
    files = []
    for v in aligned_vendors:
        dfv = subset[subset["Aligned Vendor (SUVC)"] == v].copy()
        if dfv.empty: 
            continue

        # Remove vendor's own number/name from their leads file
        cols_to_drop = ["True Vendor Number", "True Vendor Name"]
        dfv = dfv.drop(columns=cols_to_drop, errors='ignore')
        
        # Add vendor name to filename if available
        vendor_name = dfv["Aligned Supplier Name"].iloc[0] if "Aligned Supplier Name" in dfv.columns else "Unknown"
        vendor_name = str(vendor_name)  # Convert to string (handles int vendor IDs)
        # Clean name for filename
        safe_name = "".join(c for c in vendor_name if c.isalnum() or c in (' ', '-', '_')).strip()[:30]
        
        fname = os.path.join(vendor_dir, f"Vendor_{v}_{safe_name}_TRS_leads.csv")
        dfv.to_csv(fname, index=False)
        files.append(fname)
        print(f"    ✓ Vendor {v} ({safe_name}): {len(dfv)} TRS leads")
    
    return files

def build_summary(status, current_week):
    # Totals
    tot = status["Pounds_CY"].sum()

    # Aligned lbs
    iv = status["ItemVendorAligned_Lbs"].sum()       # both item+vendor
    item_any = status["ItemAligned_Lbs"].sum()       # includes 'both'
    vendor_any = status["VendorAligned_Lbs"].sum()   # includes 'both'

    # Breakouts (non-overlapping)
    i_only = max(item_any - iv, 0)
    v_only = max(vendor_any - iv, 0)
    neither = max(tot - (iv + i_only + v_only), 0)

    # KPI row with added “Any” stats
    # Count UNIQUE customers (not sum of column)
    if "Company Customer Number" in status.columns:
        # Group by customer and sum their pounds
        customer_cy_totals = status.groupby("Company Customer Number")["Pounds_CY"].sum()
        customers_cy = (customer_cy_totals > 0).sum()
        
        customer_py_totals = status.groupby("Company Customer Number")["Pounds_PY"].sum()
        customers_py = (customer_py_totals > 0).sum()
    else:
        customers_cy = 0
        customers_py = 0

    kpi = pd.DataFrame([{
        "KPI": "CY Pounds",
        "Value": tot,
        "% Any Item aligned": pct(item_any, tot),
        "% Any Vendor aligned": pct(vendor_any, tot),
        "% Item+Vendor aligned": pct(iv, tot),
        "% Item-only": pct(i_only, tot),
        "% Vendor-only": pct(v_only, tot),
        "% Neither": pct(neither, tot),
        "Current Fiscal Week": current_week,
        "Customers CY": int(customers_cy),  # ← FIXED!
        "Customers PY": int(customers_py),  # ← FIXED!
        "Customer Retention %": pct(customers_cy, customers_py),  # ← FIXED!
    }])

    # Biggest negative YoY by Company (as % of total loss)
    by_co = status.groupby(["Company Number","Company Name"], dropna=False).agg(
        CY=("Pounds_CY","sum"), PY=("Pounds_PY","sum"), d=("Delta_YoY_Lbs","sum")
    ).reset_index()
    total_loss = abs(by_co.loc[by_co["d"] < 0, "d"].sum())
    by_co["Loss_%_of_total"] = by_co.apply(
        lambda r: abs(r["d"]) / total_loss if total_loss > 0 and r["d"] < 0 else 0, axis=1
    )
    lag_sites = by_co.sort_values(["d"]).head(10)

    # Biggest customer losses
    by_cust = status.groupby(["Company Name","Customer Name"], dropna=False).agg(
        CY=("Pounds_CY","sum"), PY=("Pounds_PY","sum"), d=("Delta_YoY_Lbs","sum")
    ).reset_index().sort_values("d").head(10)

    # Narrative
    recs = []
    if total_loss > 0:
        worst = lag_sites.head(3)
        sites_list = ", ".join([f"{r['Company Name']} ({r['d']:.0f} lbs)" for _, r in worst.iterrows()]) if not worst.empty else ""
        recs.append(f"Focus on sites with largest YoY losses: {sites_list}.")
    if tot > 0:
        iv_pct = pct(iv, tot) or 0
        if iv_pct < 0.8:
            recs.append("Conversion <80% aligned on item+vendor: prioritize 'Needs Both' and 'Needs Vendor' leads.")
    summary_text = " ".join(recs) if recs else "Conversion stable; monitor 4-week momentum."

    return kpi, lag_sites, by_cust, summary_text


def build_company_yoy(status_df):
    need = ["Company Name","Pounds_CY","Pounds_PY","Delta_YoY_Lbs"]
    for c in need:
        if c not in status_df.columns:
            status_df[c] = 0
    by_co = status_df.groupby("Company Name", dropna=False).agg(
        CY_YTD=("Pounds_CY","sum"),
        PY_YTD=("Pounds_PY","sum"),
        Delta_YoY_Lbs=("Delta_YoY_Lbs","sum"),
    ).reset_index()
    by_co["YoY_Pct"] = np.where(by_co["PY_YTD"]>0, by_co["Delta_YoY_Lbs"]/by_co["PY_YTD"], np.nan)
    total_loss = abs(by_co.loc[by_co["Delta_YoY_Lbs"]<0,"Delta_YoY_Lbs"].sum())
    by_co["Loss_%_of_TotalLoss"] = np.where(
        (by_co["Delta_YoY_Lbs"]<0) & (total_loss>0),
        abs(by_co["Delta_YoY_Lbs"])/total_loss,
        0.0
    )
    by_co = by_co.sort_values(["Delta_YoY_Lbs","CY_YTD"], ascending=[True,False])
    return by_co

def build_company_weekly(df_raw):
    if "Fiscal Week Number" not in df_raw.columns:
        raise SystemExit("Missing 'Fiscal Week Number' in source.")
    current_week = int(pd.to_numeric(df_raw["Fiscal Week Number"], errors="coerce").fillna(0).max())
    tot = df_raw.groupby("Fiscal Week Number", dropna=False)[["Pounds CY","Pounds PY"]].sum().reset_index()
    tot = tot.rename(columns={"Pounds CY": "CY", "Pounds PY":"PY"})
    # NEW: add weekly Delta and YoY %
    tot["Delta_YoY_Lbs"] = tot["CY"] - tot["PY"]
    tot["YoY_Pct"] = np.where(tot["PY"]>0, tot["Delta_YoY_Lbs"]/tot["PY"], np.nan)
    return tot, current_week

def make_12w_forecast(weekly_total, current_week, method="runrate"):
    s = weekly_total.copy().sort_values("Fiscal Week Number")
    if s.empty:
        s["Type"] = "Actual"
        return s
    last4 = s.tail(4)["CY"].mean() if len(s)>=4 else s["CY"].mean()
    if len(s) >= 2:
        tail = s.tail(min(8, len(s)))
        x = tail["Fiscal Week Number"].astype(float).values
        y = tail["CY"].astype(float).values
        slope, intercept = np.polyfit(x, y, 1)
    else:
        slope, intercept = 0.0, float(s["CY"].iloc[-1])
    horizon = list(range(current_week+1, current_week+12+1))
    yhat = [slope*w + intercept for w in horizon] if method=="linear" else [last4 for _ in horizon]
    f = pd.DataFrame({"Fiscal Week Number": horizon, "CY": yhat})
    f["PY"] = np.nan
    f["Delta_YoY_Lbs"] = np.nan
    f["YoY_Pct"] = np.nan
    s["Type"] = "Actual"
    f["Type"] = "Forecast"
    return pd.concat([s, f], ignore_index=True)

# ---------------------------- NEW: Tab builders ----------------------------
def _brand_split(df):
    # Sysco vs Non-Sysco using 'Sysco Brand Indicator' (Y/N)
    if "Sysco Brand Indicator" not in df.columns:
        return pd.DataFrame(columns=["Sysco Brand Indicator","Pounds_CY","Pounds_PY","Delta_YoY_Lbs","YoY_Pct"])
    g = df.groupby("Sysco Brand Indicator", dropna=False).agg(
        Pounds_CY=("Pounds_CY","sum"),
        Pounds_PY=("Pounds_PY","sum"),
        Delta_YoY_Lbs=("Delta_YoY_Lbs","sum")
    ).reset_index()
    g["YoY_Pct"] = np.where(g["Pounds_PY"]>0, g["Delta_YoY_Lbs"]/g["Pounds_PY"], np.nan)
    return g

def _sites_rank(df):
    # Ascending by Delta YoY Pounds
    have = [c for c in ["Company Name","Pounds_CY","Pounds_PY","Delta_YoY_Lbs","YoY_Pct"] if c in df.columns]
    if "Company Name" not in have:
        return pd.DataFrame(columns=["Company Name","Pounds_CY","Pounds_PY","Delta_YoY_Lbs","YoY_Pct"])
    g = df.groupby("Company Name", dropna=False).agg(
        Pounds_CY=("Pounds_CY","sum"),
        Pounds_PY=("Pounds_PY","sum"),
        Delta_YoY_Lbs=("Delta_YoY_Lbs","sum"),
        Customers_CY=("Distinct_Customers_CY","sum"),  # NEW
        Customers_PY=("Distinct_Customers_PY","sum"),  # NEW
    ).reset_index()
    g["YoY_Pct"] = np.where(g["Pounds_PY"]>0, g["Delta_YoY_Lbs"]/g["Pounds_PY"], np.nan)
    g["Delta_Customers"] = g["Customers_CY"] - g["Customers_PY"]  # NEW
    return g.sort_values("Delta_YoY_Lbs", ascending=True)

def _items_rank(df):
    # prefer to group by Item Number + Item Description + Brand ID when available
    keys = []
    for k in ["Item Number", "Item Description", "Brand ID", "Brand"]:
        if k in df.columns:
            keys.append(k)
    if not keys:
        return pd.DataFrame(columns=["Item","Item Description","Brand ID","Pounds_CY","Pounds_PY","Delta_YoY_Lbs","YoY_Pct"])

    g = df.groupby(keys, dropna=False).agg(
        Pounds_CY=("Pounds_CY","sum"),
        Pounds_PY=("Pounds_PY","sum"),
        Delta_YoY_Lbs=("Delta_YoY_Lbs","sum")
    ).reset_index()
    g["YoY_Pct"] = np.where(g["Pounds_PY"]>0, g["Delta_YoY_Lbs"]/g["Pounds_PY"], np.nan)

    # nice column names (keep whatever exists)
    if "Item Number" in g.columns:
        g.rename(columns={"Item Number":"Item"}, inplace=True)

    return g.sort_values("Delta_YoY_Lbs", ascending=True)

def _cuisine_customer_analysis(df):
    """Analyze distinct customers by NPD Cuisine Type (TRS only)."""
    
    # Check if NPD Cuisine Type column exists
    if "NPD Cuisine Type" not in df.columns:
        return pd.DataFrame({"Note": ["NPD Cuisine Type column not found in data"]})
    
    # Group by cuisine type and count distinct customers
    cuisine_customers = []
    
    for cuisine in df["NPD Cuisine Type"].dropna().unique():
        cuisine_df = df[df["NPD Cuisine Type"] == cuisine]
        
        # Count distinct customers with pounds > 0 in each year
        customer_cy_totals = cuisine_df.groupby("Company Customer Number")["Pounds CY"].sum()
        customers_cy = (customer_cy_totals > 0).sum()
        
        customer_py_totals = cuisine_df.groupby("Company Customer Number")["Pounds PY"].sum()
        customers_py = (customer_py_totals > 0).sum()
        
        delta = customers_cy - customers_py
        pct_chg = (delta / customers_py) if customers_py > 0 else np.nan
        
        cuisine_customers.append({
            "NPD_Cuisine_Type": cuisine,
            "Customers_CY": int(customers_cy),
            "Customers_PY": int(customers_py),
            "Delta_Customers": int(delta),
            "Pct_Change": pct_chg
        })
    
    result = pd.DataFrame(cuisine_customers)
    result = result.sort_values("Delta_Customers", ascending=True)
    
    return result
    
    # Sort by biggest losses first
    result = result.sort_values("Delta_Customers", ascending=True)
    
    return result

def _overall_yoy(df):
    # DEBUG
    print("\n=== DEBUG _overall_yoy ===")
    print(f"Total rows in df: {len(df)}")
    print(f"Columns in df: {df.columns.tolist()}")
    
    row = {
        "Pounds CY": df["Pounds_CY"].sum(),
        "Pounds PY": df["Pounds_PY"].sum(),
        "Delta Pounds YoY": df["Delta_YoY_Lbs"].sum()
    }
    row["YoY %"] = (row["Delta Pounds YoY"]/row["Pounds PY"]) if row["Pounds PY"] else np.nan
    
    # Recalculate distinct customers from raw customer numbers
    # (can't sum Distinct_Customers_CY because customers have multiple item rows)
    if "Company Customer Number" in df.columns:
        # Count unique customers who bought in CY (sum their pounds across all items > 0)
        customer_cy_totals = df.groupby("Company Customer Number")["Pounds_CY"].sum()
        customers_cy = (customer_cy_totals > 0).sum()
        
        # Count unique customers who bought in PY
        customer_py_totals = df.groupby("Company Customer Number")["Pounds_PY"].sum()
        customers_py = (customer_py_totals > 0).sum()
        
        print(f"Calculated Customers CY: {customers_cy}")
        print(f"Calculated Customers PY: {customers_py}")
        
        row["Customers CY"] = int(customers_cy)
        row["Customers PY"] = int(customers_py)
        row["Delta Customers"] = int(customers_cy - customers_py)
        row["Customer Retention %"] = (customers_cy / customers_py) if customers_py > 0 else np.nan
        row["Avg Lbs/Customer CY"] = (row["Pounds CY"] / customers_cy) if customers_cy > 0 else 0
        row["Avg Lbs/Customer PY"] = (row["Pounds PY"] / customers_py) if customers_py > 0 else 0
    else:
        row["Customers CY"] = 0
        row["Customers PY"] = 0
        row["Delta Customers"] = 0
        row["Customer Retention %"] = np.nan
        row["Avg Lbs/Customer CY"] = 0
        row["Avg Lbs/Customer PY"] = 0
    
    print("=========================\n")
    return pd.DataFrame([row])

def _write_account_tab(xw, tab_name, df, filtered_company=None):
    start = 0
    # 1) OVERALL
    overall = _overall_yoy(df)
    overall.to_excel(xw, sheet_name=tab_name, index=False, startrow=start)
    _format_table_at(
        xw.book[tab_name], header_row_0idx=start, n_rows=overall.shape[0],
        number_headers={"Pounds CY","Pounds PY","Delta Pounds YoY","Customers CY","Customers PY","Delta Customers","Avg Lbs/Customer CY","Avg Lbs/Customer PY"},
        percent_headers={"YoY %","Customer Retention %"}
    )
    start += overall.shape[0] + 2

    # 2) BRAND SPLIT
    brand = _brand_split(df)
    if not brand.empty:
        brand.to_excel(xw, sheet_name=tab_name, index=False, startrow=start)
        _format_table_at(
            xw.book[tab_name], header_row_0idx=start, n_rows=brand.shape[0],
            number_headers={"Pounds_CY","Pounds_PY","Delta_YoY_Lbs"},
            percent_headers={"YoY_Pct"}
        )
        start += brand.shape[0] + 2

    # 3) SITES
    sites = _sites_rank(df)
    if not sites.empty:
        sites.to_excel(xw, sheet_name=tab_name, index=False, startrow=start)
        _format_table_at(
            xw.book[tab_name], header_row_0idx=start, n_rows=sites.shape[0],
            number_headers={"Pounds_CY","Pounds_PY","Delta_YoY_Lbs","Customers_CY","Customers_PY","Delta_Customers"},
            percent_headers={"YoY_Pct"}
        )
        start += sites.shape[0] + 2

    # 4) ITEMS
    items = _items_rank(df)
    if not items.empty:
        items.to_excel(xw, sheet_name=tab_name, index=False, startrow=start)
        _format_table_at(
            xw.book[tab_name], header_row_0idx=start, n_rows=items.shape[0],
            number_headers={"Pounds_CY","Pounds_PY","Delta_YoY_Lbs"},
            percent_headers={"YoY_Pct"}
        )
        
def _write_weekly_chart_only_tab(xw, sheet_name, weekly_df, forecast_method_desc, filtered_company=None):
    from openpyxl.chart import BarChart, LineChart, Reference
    from openpyxl.chart.shapes import GraphicalProperties
    
    # Write data at top
    df_out = weekly_df.copy()
    df_out.to_excel(xw, sheet_name=sheet_name, index=False, startrow=0)
    
    ws = xw.book[sheet_name]
    
    # ADD SITE FILTER NOTE AT TOP
    if filtered_company:
        ws.cell(1, 15, value=f"FILTERED TO: {filtered_company}").font = Font(bold=True, size=14, color="C00000")
    
    # Find columns
    headers = {cell.value: cell.col_idx for cell in ws[1] if cell.value is not None}
    col_week = headers.get("Fiscal Week Number")
    col_cy = headers.get("CY")
    col_delta = headers.get("Delta_YoY_Lbs")
    col_yoy = headers.get("YoY_Pct")
    col_type = headers.get("Type")
    
    if not all([col_week, col_delta, col_yoy]):
        print(f"      ⚠ Missing columns for chart in {sheet_name}")
        return
    
    n_rows = df_out.shape[0]
    actual_count = df_out[df_out['Type'] == 'Actual'].shape[0] if col_type else n_rows
    
    # === CREATE HELPER COLUMNS FOR POSITIVE/NEGATIVE BARS ===
    ws.cell(1, 12, value="▲ Growth (Green Bars)")
    ws.cell(1, 13, value="▼ Decline (Red Outlined Bars)")
    
    for i in range(2, n_rows+2):
        val = ws.cell(i, col_delta).value
        if val and float(val) >= 0:
            ws.cell(i, 12, value=val)
            ws.cell(i, 13, value=None)
        else:
            ws.cell(i, 12, value=None)
            ws.cell(i, 13, value=val)
    
    # === BAR CHART ===
    bar = BarChart()
    bar.type = "col"
    bar.grouping = "clustered"
        
    bar.width = 30
    bar.height = 18
    bar.gapWidth = 30
    
    # POSITIVE bars (green)
    pos_data = Reference(ws, min_col=12, min_row=1, max_row=n_rows+1)
    bar.add_data(pos_data, titles_from_data=True)
    bar.series[0].graphicalProperties.solidFill = "92D050"
    bar.series[0].graphicalProperties.line.solidFill = "92D050"
    bar.series[0].graphicalProperties.line.width = 12700

    # NEGATIVE bars (red) - SIMPLIFIED FIX
    neg_data = Reference(ws, min_col=13, min_row=1, max_row=n_rows+1)
    bar.add_data(neg_data, titles_from_data=True)
    bar.series[1].graphicalProperties.solidFill = "FF0000"  # Pure red
    bar.series[1].graphicalProperties.line.solidFill = "FF0000"
    bar.series[1].graphicalProperties.line.width = 12700

    # Ensure no transparency
    if hasattr(bar.series[1].graphicalProperties, 'noFill'):
        bar.series[1].graphicalProperties.noFill = False
    
    # X-axis
    x_ref = Reference(ws, min_col=col_week, min_row=2, max_row=n_rows+1)
    bar.set_categories(x_ref)
    
    # === PRIMARY Y-AXIS (SHORT LABEL) ===
    bar.y_axis.title = "Lbs."
    bar.y_axis.number_format = '#,##0'
    bar.y_axis.delete = False
    
    from openpyxl.chart.axis import ChartLines
    bar.y_axis.majorGridlines = ChartLines()
    
    # === X-AXIS (SHORT LABEL) ===
    bar.x_axis.title = "FWk"
    bar.x_axis.delete = False
    bar.x_axis.tickLblPos = "low"
    
    # === ACTUAL CY LINE ===
    actual_line = LineChart()
    actual_line.y_axis.axId = 100
    
    ws.cell(1, 10, value="CY Pounds (Actual)")
    actual_data = Reference(ws, min_col=10, min_row=1, max_row=actual_count+1)
    for i in range(2, actual_count+2):
        ws.cell(i, 10, value=ws.cell(i, col_cy).value)
    
    actual_line.add_data(actual_data, titles_from_data=True)
    actual_line.set_categories(Reference(ws, min_col=col_week, min_row=2, max_row=actual_count+1))
    
    try:
        actual_line.series[0].graphicalProperties.line.width = 30000
        actual_line.series[0].graphicalProperties.line.solidFill = "4472C4"
    except:
        pass
    
    bar += actual_line
    
    # === FORECAST LINE ===
    if actual_count < n_rows:
        forecast_line = LineChart()
        forecast_line.y_axis.axId = 100
        
        ws.cell(1, 11, value="CY Pounds (Forecast)")
        forecast_data = Reference(ws, min_col=11, min_row=1, max_row=n_rows+1)
        
        ws.cell(actual_count+1, 11, value=ws.cell(actual_count+1, col_cy).value)
        for i in range(actual_count+2, n_rows+2):
            ws.cell(i, 11, value=ws.cell(i, col_cy).value)
        
        forecast_line.add_data(forecast_data, titles_from_data=True)
        forecast_line.set_categories(Reference(ws, min_col=col_week, min_row=actual_count+1, max_row=n_rows+1))
        
        try:
            forecast_line.series[0].graphicalProperties.line.width = 35000
            forecast_line.series[0].graphicalProperties.line.dashStyle = "dash"
            forecast_line.series[0].graphicalProperties.line.solidFill = "ED7D31"
        except:
            pass
        
        bar += forecast_line
    
    # === YoY % LINE (SHORT LABEL) ===
    yoy_line = LineChart()
    yoy_line.y_axis.axId = 200
    yoy_line.y_axis.crosses = "max"
    yoy_line.y_axis.title = "% Chg"
    yoy_line.y_axis.number_format = '0%'
    yoy_line.y_axis.delete = False
    yoy_line.y_axis.majorGridlines = ChartLines()
    
    yoy_data = Reference(ws, min_col=col_yoy, min_row=1, max_row=n_rows+1)
    yoy_line.add_data(yoy_data, titles_from_data=True)
    yoy_line.set_categories(x_ref)
    
    try:
        yoy_line.series[0].graphicalProperties.line.width = 25000
        yoy_line.series[0].graphicalProperties.line.solidFill = "70AD47"
    except:
        pass
    
    bar += yoy_line
    
    # === LEGEND ===
    bar.legend.position = 'tr'
    try:
        bar.legend.spPr = GraphicalProperties(solidFill="FFFFFF")
    except:
        pass
    
    ws.add_chart(bar, f"A{n_rows + 4}")
    
    # === EXPLANATION TEXT BOX ===
    explain_row = n_rows + 38
    
    # Header
    cell = ws.cell(explain_row, 1, value="HOW TO READ THIS CHART:")
    cell.font = Font(bold=True, size=12, color="000000")
    
    # Left axis explanation
    cell = ws.cell(explain_row + 1, 1, value="Left Axis (Pounds): Shows both volume changes compared to last year AND current year actual pounds")
    cell.font = Font(size=10)
    
    # Right axis explanation
    cell = ws.cell(explain_row + 2, 1, value="Right Axis (Growth %): Shows year-over-year percentage change")
    cell.font = Font(size=10)
    
    # Color key
    cell = ws.cell(explain_row + 3, 1, value="Green Bars = Growth vs Last Year  |  Red Bars = Decline vs Last Year  |  Blue Line = This Year Actual  |  Orange Dashed = Forecast  |  Green Line = % Growth")
    cell.font = Font(size=10, bold=True)
    
    # Data range
    cell = ws.cell(explain_row + 5, 1, value=f"Actual Data: Weeks 1-{actual_count}  |  Forecast: Weeks {actual_count+1}-{n_rows}")
    cell.font = Font(size=10, italic=True)
    
    if filtered_company:
        cell = ws.cell(explain_row + 6, 1, value=f"⚠ DATA FILTERED TO: {filtered_company}")
        cell.font = Font(size=11, bold=True, color="C00000")
    
    print(f"      ✓ Chart with color-coded bars and user-friendly labels")
        
def _format_table_at(ws, header_row_0idx: int, n_rows: int,
                     number_headers=None, percent_headers=None):
    from openpyxl.styles import Alignment  # Add this import
    
    number_headers = set(number_headers or [])
    percent_headers = set(percent_headers or [])
    header_row_1idx = header_row_0idx + 1
    first_data = header_row_1idx + 1
    last_data  = header_row_1idx + n_rows
    if n_rows <= 0:
        return

    # Map header text in THIS table only
    header_to_col = {str(c.value).strip(): c.col_idx
                     for c in ws[header_row_1idx] if c.value is not None}

    def _coerce(cell, as_percent: bool):
        v = cell.value
        if v is None or isinstance(v, (int, float)):
            return
        s = str(v).replace(",", "").strip()
        if not s:
            return
        try:
            if as_percent and s.endswith("%"):
                cell.value = float(s[:-1]) / 100.0
            else:
                cell.value = float(s)
        except Exception:
            pass

    # Numbers
    for h in number_headers:
        col = header_to_col.get(h)
        if not col:
            continue
        for row in ws.iter_rows(min_row=first_data, max_row=last_data, min_col=col, max_col=col):
            _coerce(row[0], as_percent=False)
            row[0].number_format = "#,##0"

    # Percents
    for h in percent_headers:
        col = header_to_col.get(h)
        if not col:
            continue
        for row in ws.iter_rows(min_row=first_data, max_row=last_data, min_col=col, max_col=col):
            _coerce(row[0], as_percent=True)
            row[0].number_format = "0.0%"
    
    # NEW: Text wrapping for all cells in this table
    for row in ws.iter_rows(min_row=header_row_1idx, max_row=last_data):
        for cell in row:
            if cell.value and isinstance(cell.value, str) and len(str(cell.value)) > 30:
                cell.alignment = Alignment(wrap_text=True, vertical='top')
def create_guide_tab(xw, status_df, weeks_covered, filtered_company_name=None, min_week=1, max_week=52):
    """
    Create an instructional guide tab that teaches category managers how to use the report.
    This should be the FIRST tab in the workbook.
    """
    from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    
    # Create or get the Guide sheet
    if "Guide" in xw.book.sheetnames:
        ws = xw.book["Guide"]
    else:
        ws = xw.book.create_sheet("Guide", 0)  # Insert as first sheet
    
    # Styling
    title_font = Font(name="Calibri", size=18, bold=True, color="0066CC")
    header_font = Font(name="Calibri", size=14, bold=True, color="FFFFFF")
    header_fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")
    subheader_font = Font(name="Calibri", size=12, bold=True, color="333333")
    body_font = Font(name="Calibri", size=11)
    italic_font = Font(name="Calibri", size=10, italic=True, color="666666")

    def add_section_header(ws, row_num, text):
        """Add a formatted section header with proper spacing."""
        cell = ws.cell(row_num, 1, value=text)
        cell.font = header_font
        cell.fill = header_fill
        cell.alignment = Alignment(horizontal='left', vertical='center', wrap_text=False)
        ws.merge_cells(f"A{row_num}:G{row_num}")
        ws.row_dimensions[row_num].height = 30  # Consistent header height
        return row_num + 1
    
    # === REPORT HEADER ===
    row = 1
    title_cell = ws.cell(row, 1,"📊 HOW TO USE THIS CATEGORY PERFORMANCE REPORT")
    title_cell.font = title_font
    title_cell.alignment = Alignment(horizontal='center', vertical='center', wrap_text=False)
    ws.merge_cells(f"A{row}:G{row}")
    ws.row_dimensions[row].height = 50  # Extra tall for impact
    row += 1
    
    ws.cell(row, 1, value=f"Generated: {datetime.now():%B %d, %Y at %I:%M %p}").font = italic_font
    ws.row_dimensions[row].height = 20
    row += 1
    
    if filtered_company_name:
        ws.cell(row, 1, value=f"🎯 FILTERED TO: {filtered_company_name}").font = Font(bold=True, size=12, color="C00000")
        row += 1
    
    ws.cell(row, 1, value=f"📅 Coverage: Weeks {min_week}-{max_week} ({weeks_covered} weeks of data)").font = body_font
    row += 2
    
    # === SECTION 1: QUICK START ===
    row = add_section_header(ws, row, "🚀 QUICK START: Your Monday Morning Action Plan")
 
    # Generate top 3 actions based on actual data
    actions = []
    
    # Action 1: Biggest declining customer who's still active
    declining = status_df[status_df["Delta_YoY_Lbs"] < 0].copy()
    if not declining.empty:
        declining = declining.sort_values("Delta_YoY_Lbs").head(1)
        cust_name = declining.iloc[0]["Customer Name"]
        cust_loss = abs(declining.iloc[0]["Delta_YoY_Lbs"])
        dsm = declining.iloc[0].get("Customer DSM", "your DSM")
        actions.append(f"Win Back '{cust_name}' – Down {cust_loss:,.0f} lbs YoY but still active. Call {dsm} TODAY.")
    
    # Action 2: Conversion opportunity
    needs_both = status_df[status_df["Conversion_Status"] == "Needs Both"]
    if not needs_both.empty:
        opp_lbs = needs_both["Pounds_CY"].sum()
        actions.append(f"Convert {len(needs_both):,} accounts to aligned items+vendors ({opp_lbs:,.0f} lbs opportunity).")
    
    # Action 3: Focus on biggest declining item/vendor
    if "Item Number" in status_df.columns:
        items_declining = status_df.groupby(["Item Number", "Item Description"]).agg(
            Delta=("Delta_YoY_Lbs", "sum")
        ).reset_index().sort_values("Delta")
        if not items_declining.empty and items_declining.iloc[0]["Delta"] < 0:
            top_item = items_declining.iloc[0]["Item Description"]
            item_loss = abs(items_declining.iloc[0]["Delta"])
            actions.append(f"Investigate '{top_item}' – Down {item_loss:,.0f} lbs. Check pricing/availability/competition.")
    
    for i, action in enumerate(actions[:3], 1):
        ws.cell(row, 1, value=f"   {i}. {action}").font = Font(size=11, bold=True)
        ws.cell(row, 1).alignment = Alignment(wrap_text=True, vertical="top")
        row += 1
    
    row += 1
    
    # === SECTION 2: UNDERSTANDING THE METRICS ===
    row = add_section_header(ws, row, "📖 KEY METRICS EXPLAINED")
    metrics_guide = [
        ("Pounds CY / PY", "Current Year vs Prior Year volume in pounds"),
        ("Delta YoY Lbs", "Change from last year (negative = declining, positive = growing)"),
        ("YoY %", "Percentage change from last year"),
        ("Conversion Status", "'Converted' = using aligned item+vendor | 'Needs Both' = top priority to convert"),
        ("% of Expected", "(Award tab) Your progress vs. pro-rated goal for this point in the year"),
        ("Last Invoice Date", "Most recent purchase of ANY category - shows if customer is still active with Sysco"),
        ("Active Customer", "Last Invoice Date within past 8 weeks = still buying from us (even if not this category)"),
    ]
    
    for metric, explanation in metrics_guide:
        ws.cell(row, 1, value=f"   • {metric}:").font = Font(bold=True, size=10)
        ws.cell(row, 2, value=explanation).font = body_font
        ws.merge_cells(f"B{row}:G{row}")
        row += 1
    
    row += 1
    
   # === SECTION 3: TAB-BY-TAB GUIDE ===
    row = add_section_header(ws, row, "📑 HOW TO USE EACH TAB")

    tab_guide = [
        ("Guide (this tab)", "Read this first. Explains metrics, filters, and how to use the report."),
        ("Summary", "Overall health snapshot: YoY totals, Sysco brand split, biggest declining sites/customers."),
        ("01_All_Accounts", "All account types combined. Overall metrics, brand split, site rankings, item rankings."),
        ("02_TRS", "TRS accounts only. Same metrics as tab 01 but filtered to operators."),
        ("03_LCC", "LCC accounts only. Local commercial customers."),
        ("04_CMU", "CMU accounts only. Street/street-adjacent accounts."),
        ("05-09 Weekly Charts", "Weekly trends + 12-week forecast. One chart per account type (All/TRS/CMU/LCC). Green bars = growth, red = decline."),
        ("10_VendorItem_Alignment", "Full sales leads with alignment status. Use this to identify conversion opportunities."),
        ("11_Award_vs_Sales", "Track vendor commitments. Use '% of Expected' for partial-year reports (pro-rated)."),
        ("12_Sites_by_Leads", "Sites ranked by lead volume. Shows which sites have the most opportunities."),
        ("13_DSM_Opportunity", "TRS-only: Win-back + conversion targets by DSM. Top 5 customers per DSM listed."),
        ("14_Vendor_Leads_Index", "List of vendor-specific CSV files created (TRS only). Send to vendors for their leads."),
    ]

    for tab_name, guidance in tab_guide:
        ws.cell(row, 1, value=f"   {tab_name}:").font = Font(bold=True, size=10, color="0066CC")
        ws.cell(row, 2, value=guidance).font = body_font
        ws.cell(row, 2).alignment = Alignment(wrap_text=True)
        ws.merge_cells(f"B{row}:G{row}")
        row += 1

    row += 1
    
    # === SECTION 4: WHAT ARE SALES LEADS ===
    row = add_section_header(ws, row, "📋 SALES LEADS TAB: What's Included")
    ws.cell(row, 1, value="   INCLUDED:").font = Font(bold=True, size=11, color="008000")
    row += 1
    ws.cell(row, 1, value="      • Customers currently buying this category (minimum 20 lbs/week)").font = body_font
    row += 1
    ws.cell(row, 1, value="      • TRS + LCC accounts only (operators you're allowed to contact)").font = body_font
    row += 1
    ws.cell(row, 1, value="      • Conversion Status: 'Needs Both', 'Needs Item', 'Needs Vendor', or 'Partial'").font = body_font
    row += 1
    ws.cell(row, 1, value="      • All sites (unless report was run with --company filter)").font = body_font
    row += 2

    ws.cell(row, 1, value="   EXCLUDED:").font = Font(bold=True, size=11, color="C00000")
    row += 1
    ws.cell(row, 1, value="      • CMU accounts (not street accounts - can't be contacted directly)").font = body_font
    row += 1
    ws.cell(row, 1, value="      • Customers with zero current volume or already 'Converted' (100% aligned)").font = body_font
    row += 1
    ws.cell(row, 1, value="      • Customers below 20 lbs/week threshold").font = body_font
    row += 2

    ws.cell(row, 1, value="   DSM Opportunity Tab (separate): Shows TRS win-back targets + conversion opps by DSM").font = Font(italic=True, size=10, color="666666")
    row += 2

    # === SECTION 5: THE GROWTH FORMULA (CORRECTED) ===
    row = add_section_header(ws, row, "📈 THE CATEGORY GROWTH FORMULA (In Priority Order)")    
    ws.cell(row, 1, value="   1. WIN BACK declining customers who are still ACTIVE (HIGHEST PRIORITY)").font = Font(bold=True, size=11, color="C00000")
    row += 1
    ws.cell(row, 1, value="      → Focus on customers down in THIS category but Last Invoice Date shows recent purchases").font = body_font
    row += 1
    ws.cell(row, 1, value="      → They're still buying from Sysco (just not YOUR category) - easiest to win back").font = body_font
    row += 1
    ws.cell(row, 1, value="      → Target customers representing large % of total category losses (see Summary tab)").font = body_font
    row += 1
    ws.cell(row, 1, value="      → Send to Site within 48 hours or Send and Meet/Call - these are HOT LEADS").font = body_font
    row += 2
    
    ws.cell(row, 1, value="   2. CONVERT to aligned items/vendors (CLOSE SECOND PRIORITY)").font = Font(bold=True, size=11, color="E67E22")
    row += 1
    ws.cell(row, 1, value="      → Use Sales Leads tab. Target 'Needs Both' first (neither item nor vendor aligned)").font = body_font
    row += 1
    ws.cell(row, 1, value="      → Better margins, better service, secures the business long-term").font = body_font
    row += 1
    ws.cell(row, 1, value="      → Vendors want these leads - route to them via Vendor Leads folder").font = body_font
    row += 2
    
    ws.cell(row, 1, value="   3. TARGET high-volume items in decline (INVESTIGATE ROOT CAUSE)").font = Font(bold=True, size=11, color="3498DB")
    row += 1
    ws.cell(row, 1, value="      → See Items tab (in Account Tabs). Sort by largest declines").font = body_font
    row += 1
    ws.cell(row, 1, value="      → Check: Pricing competitive? Availability issues? New competitor?").font = body_font
    row += 1
    ws.cell(row, 1, value="      → If item is losing across multiple sites = systemic issue to escalate").font = body_font
    row += 2
    
    ws.cell(row, 1, value="   ⚠️ WHY 'ACTIVE' MATTERS:").font = Font(bold=True, size=11, color="666666")
    row += 1
    ws.cell(row, 1, value="      Customers who haven't bought ANYTHING in 8+ weeks are likely lost to a competitor.").font = body_font
    row += 1
    ws.cell(row, 1, value="      Focus your time on active customers - they're still reachable and most likely to convert.").font = body_font
    row += 2
    
    # === SECTION 6: FAQs (CORRECTED) ===
    row = add_section_header(ws, row, "❓ FREQUENTLY ASKED QUESTIONS")    
    faqs = [
        ("Q: How do I read the Summary tab?", 
         "A: Start with Overall YoY metrics (top section). Then check Sysco Brand split to see if losses are brand-specific. Bottom tables show biggest losing sites and customers - prioritize these."),
        
        ("Q: Why does Award vs Sales show low %?", 
         "A: Check report coverage at top of tab. For partial-year reports, use '% of Expected' column (pro-rated), not '% of Annual Award'. Green = on track, Red = under 80%."),
        
        ("Q: What does 'Needs Both' mean?", 
         "A: Customer bought this category but used NEITHER the aligned item NOR aligned vendor. Highest conversion priority - most margin opportunity."),
        
        ("Q: How do I know if a customer is still active?", 
         "A: Check 'Last Invoice Date' in Sales Leads tab. If within past 8 weeks, they're active (buying other categories even if declining in yours)."),
        
        ("Q: Who gets the Sales Leads?", 
         "A: Route to Sites/DSMs based on Customer DSM and Territory columns. These are pre-filtered to active customers with meaningful volume."),
        
        ("Q: Why TRS only for vendor leads?", 
         "A: Vendors can't call on LCC/CMU street accounts. TRS = operator accounts vendors are allowed to contact directly."),
        
    ]
    
    for question, answer in faqs:
        ws.cell(row, 1, value=question).font = Font(bold=True, size=10)
        ws.merge_cells(f"A{row}:G{row}")
        row += 1
        ws.cell(row, 1, value=answer).font = body_font
        ws.cell(row, 1).alignment = Alignment(wrap_text=True)
        ws.merge_cells(f"A{row}:G{row}")
        row += 1
    
    row += 1
    
    # === SECTION 7: CONTACT INFO ===
    row = add_section_header(ws, row, "📞 NEED HELP?")
    
    ws.cell(row, 1, value="   Questions about this report? Contact your Category Manager.").font = body_font
    ws.merge_cells(f"A{row}:G{row}")
    row += 1
    
    ws.cell(row, 1, value="   📧 Report generated automatically. Do not reply to notification emails.").font = italic_font
    ws.merge_cells(f"A{row}:G{row}")
    

    # === FINAL FORMATTING ===
    
    # Apply text wrapping to all content cells
    for row_cells in ws.iter_rows(min_row=1, max_row=ws.max_row):
        for cell in row_cells:
            if cell.value and not cell.fill.start_color.rgb == "4472C4":  # Not a header
                cell.alignment = Alignment(wrap_text=True, vertical='top')
    
    # Set minimum heights for content rows (not headers)
    for row_num in range(1, ws.max_row + 1):
        cell = ws.cell(row_num, 1)
        
        # Skip if it's a blue header (already sized)
        if cell.fill and hasattr(cell.fill, 'start_color') and cell.fill.start_color.rgb == "4472C4":
            continue
        
        # Calculate based on content length
        if cell.value:
            text_length = len(str(cell.value))
            if text_length > 100:  # Long text needs more space
                ws.row_dimensions[row_num].height = max(30, min(text_length / 4, 80))
            elif text_length > 50:
                ws.row_dimensions[row_num].height = 25
            else:
                ws.row_dimensions[row_num].height = 18
    
    print("  ✓ Guide tab created")

def build_dsm_opportunity_scorecard(status_df, raw_df, source_path, active_weeks=8, min_ytd=260):
    """Build COMBINED DSM scorecard..."""
    
    # === CRITICAL: Filter to TRS ONLY at the very start ===
    if "Customer Account Type Code" not in status_df.columns:
        print(f"    ❌ Customer Account Type Code not in status_df - aborting DSM scorecard")
        return (pd.DataFrame(), pd.DataFrame(), pd.DataFrame(), pd.DataFrame())
    
    # Filter BOTH status_df and raw_df to TRS only
    status_df = status_df[status_df["Customer Account Type Code"] == "TRS"].copy()
    raw_df = raw_df[raw_df["Customer Account Type Code"] == "TRS"].copy()
    
    print(f"    ✓ Filtered to TRS only: {len(status_df):,} status records, {len(raw_df):,} raw records")
    
    if status_df.empty:
        print(f"    ℹ No TRS accounts found")
        return (pd.DataFrame(), pd.DataFrame(), pd.DataFrame(), pd.DataFrame())
    
    # === REST OF FUNCTION CONTINUES (no changes below this) ===
    import pandas as pd
    import numpy as np
    from datetime import datetime, timedelta
    
    # Read ORIGINAL CSV directly (before any alignment merging)
    print(f"    Reading original CSV for win-back calculation...")
    truly_raw = pd.read_csv(source_path, dtype=str)
    truly_raw.columns = truly_raw.columns.str.replace(r"\s+", " ", regex=True).str.strip()

    # Convert to numeric (strip commas first)
    for col in ["Pounds CY", "Pounds PY", "Fiscal Week Number"]:
        if col in truly_raw.columns:
            truly_raw[col] = clean_numeric_column(truly_raw[col])

    # DEBUG: Check customer 010320895 before filtering
    test = truly_raw[truly_raw["Company Customer Number"].str.strip() == "010320895"]
    if len(test) > 0:
        print(f"\n    === CUSTOMER 010320895 ===")
        print(f"    Rows: {len(test)} | PY: {test['Pounds PY'].sum():.0f} | CY: {test['Pounds CY'].sum():.0f} | Delta: {test['Pounds CY'].sum() - test['Pounds PY'].sum():.0f}")
        print(f"    Account Type: {test['Customer Account Type Code'].iloc[0]} | DSM: '{test['Customer DSM'].iloc[0]}'")

    # Filter to TRS
    if "Customer Account Type Code" in truly_raw.columns:
        truly_raw = truly_raw[truly_raw["Customer Account Type Code"] == "TRS"].copy()
        print(f"    ⚠ Filtered to TRS: {len(truly_raw):,} records")

    # WIN-BACK: Use truly raw data
    cutoff_date = datetime.now() - timedelta(weeks=active_weeks)

    invoice_dates = truly_raw[["Company Number", "Company Customer Number", "Last Invoice Date"]].copy()
    invoice_dates["Last Invoice Date"] = pd.to_datetime(invoice_dates["Last Invoice Date"], errors="coerce")
    last_invoice = invoice_dates.groupby(["Company Number", "Company Customer Number"])["Last Invoice Date"].max().reset_index()

    winback_customers = truly_raw.groupby(
        ["Company Number", "Company Name", "Customer Name", "Company Customer Number",
        "Customer DSM", "Customer Territory"],
        dropna=False
    ).agg({
        "Pounds CY": "sum",
        "Pounds PY": "sum"
    }).reset_index()

    winback_customers = winback_customers.rename(columns={"Pounds CY": "Pounds_CY", "Pounds PY": "Pounds_PY"})
    winback_customers = winback_customers.merge(last_invoice, on=["Company Number", "Company Customer Number"], how="left")
    winback_customers["Delta_YoY_Lbs"] = winback_customers["Pounds_CY"] - winback_customers["Pounds_PY"]
    winback_customers["YoY_Pct"] = np.where(
        winback_customers["Pounds_PY"] > 0,
        winback_customers["Delta_YoY_Lbs"] / winback_customers["Pounds_PY"],
        np.nan
    )
    winback_customers["Is_Active"] = (winback_customers["Last Invoice Date"] >= cutoff_date).astype(int)

    winback_opps = winback_customers[
        (winback_customers["Is_Active"] == 1) &
        (winback_customers["Delta_YoY_Lbs"] < 0) &
        (winback_customers["Pounds_PY"] >= min_ytd)
    ].copy()
    winback_opps["Opportunity_Lbs"] = abs(winback_opps["Delta_YoY_Lbs"])

    print(f"    Win-back (from truly raw): {len(winback_opps)} customers")

    if not winback_opps.empty:
        winback_items = []
        for _, customer in winback_opps.head(100).iterrows():
            cust_num = customer["Company Customer Number"]
            
            # Get items this customer bought PY but dropped in CY
            cust_items = truly_raw[truly_raw["Company Customer Number"].str.strip() == cust_num].copy()
            
            if not cust_items.empty:
                item_summary = cust_items.groupby("Item Description").agg({
                    "Pounds PY": "sum",
                    "Pounds CY": "sum"
                }).reset_index()
                
                # Items they stopped buying (PY > 0, CY = 0 or very low)
                lost_items = item_summary[
                    (item_summary["Pounds PY"] > 0) & 
                    (item_summary["Pounds CY"] < item_summary["Pounds PY"] * 0.2)
                ].sort_values("Pounds PY", ascending=False).head(3)
                
                if not lost_items.empty:
                    items_text = " | ".join([
                        f"{row['Item Description'][:40]} ({row['Pounds PY']:.0f} lbs lost)"
                        for _, row in lost_items.iterrows()
                    ])
                    winback_items.append({
                        "Company Customer Number": cust_num,
                        "Items_Lost": items_text
                    })
        
        if winback_items:
            winback_items_df = pd.DataFrame(winback_items)
            winback_opps = winback_opps.merge(
                winback_items_df, 
                on="Company Customer Number", 
                how="left"
            )
            winback_opps["Items_Lost"] = winback_opps["Items_Lost"].fillna("")
    
    # CONVERSION: Use status_df (needs alignment data)  
    customer_totals = status_df.groupby(
        ["Company Number", "Company Name", "Customer Name", "Company Customer Number", 
         "Customer DSM", "Customer Territory"],
        dropna=False
    ).agg({
        "Pounds_CY": "sum",
        "Pounds_PY": "sum"
    }).reset_index()
    
    customer_totals = customer_totals.merge(last_invoice, on=["Company Number", "Company Customer Number"], how="left")
    customer_totals["Delta_YoY_Lbs"] = customer_totals["Pounds_CY"] - customer_totals["Pounds_PY"]
    customer_totals["YoY_Pct"] = np.where(
        customer_totals["Pounds_PY"] > 0,
        customer_totals["Delta_YoY_Lbs"] / customer_totals["Pounds_PY"],
        np.nan
    )
    customer_totals["Is_Active"] = (customer_totals["Last Invoice Date"] >= cutoff_date).astype(int)
    
    if "Conversion_Status" in status_df.columns:
        conv_status = status_df.groupby("Company Customer Number")["Conversion_Status"].first().reset_index()
        customer_totals = customer_totals.merge(conv_status, on="Company Customer Number", how="left")
    
    print(f"    Total unique customers: {len(customer_totals):,}")
    print(f"    Active customers: {(customer_totals['Is_Active'] == 1).sum():,}")
    print(f"    Declining customers: {(customer_totals['Delta_YoY_Lbs'] < 0).sum():,}")
    
    conversion_opps = customer_totals[
        (customer_totals["Pounds_CY"] >= min_ytd) &
        (customer_totals.get("Conversion_Status", "").isin(["Needs Both", "Needs Item", "Needs Vendor"]))
    ].copy()
    conversion_opps["Opportunity_Lbs"] = conversion_opps["Pounds_CY"]
    
    print(f"    Conversion opportunities: {len(conversion_opps):,} customers")
    
    # ADD THIS SECTION: Get item details for conversion targets
    if not conversion_opps.empty and not status_df.empty:
        conversion_items = []
        for _, customer in conversion_opps.head(100).iterrows():
            cust_num = customer["Company Customer Number"]
            
            # Get current items needing conversion
            cust_status = status_df[
                (status_df["Company Customer Number"] == cust_num) &
                (status_df["Conversion_Status"].isin(["Needs Both", "Needs Item", "Needs Vendor"]))
            ].copy()
            
            if not cust_status.empty:
                # Get top 3 items by volume
                top_items = cust_status.nlargest(3, "Pounds_CY")
                
                items_text = " | ".join([
                    f"{row.get('Item Description', 'Unknown')[:30]} → {row.get('SUPC', 'N/A')} ({row['Pounds_CY']:.0f} lbs)"
                    for _, row in top_items.iterrows()
                ])
                conversion_items.append({
                    "Company Customer Number": cust_num,
                    "Items_To_Convert": items_text
                })
        
        if conversion_items:
            conversion_items_df = pd.DataFrame(conversion_items)
            conversion_opps = conversion_opps.merge(
                conversion_items_df,
                on="Company Customer Number",
                how="left"
            )
            conversion_opps["Items_To_Convert"] = conversion_opps["Items_To_Convert"].fillna("(Item detail not available)")
    
    print(f"    Conversion opportunities: {len(conversion_opps):,} customers")
    
    if winback_opps.empty and conversion_opps.empty:
        return (pd.DataFrame(), pd.DataFrame(), pd.DataFrame(), pd.DataFrame())
    
    # DSM SUMMARY
    dsm_winback = winback_opps.groupby(["Company Name", "Customer DSM"], dropna=False).agg(
        WinBack_Lbs=("Opportunity_Lbs", "sum"),
        WinBack_Customers=("Company Customer Number", "nunique")
    ).reset_index()
    
    dsm_conversion = conversion_opps.groupby(["Company Name", "Customer DSM"], dropna=False).agg(
        Conversion_Lbs=("Opportunity_Lbs", "sum"),
        Conversion_Customers=("Company Customer Number", "nunique")
    ).reset_index()
    
    dsm_summary = dsm_winback.merge(dsm_conversion, on=["Company Name", "Customer DSM"], how="outer").fillna(0)
    dsm_summary["Total_Opportunity_Lbs"] = dsm_summary["WinBack_Lbs"] + dsm_summary["Conversion_Lbs"]
    dsm_summary["Total_Customers"] = (dsm_summary["WinBack_Customers"] + dsm_summary["Conversion_Customers"]).astype(int)
    dsm_summary = dsm_summary.sort_values("Total_Opportunity_Lbs", ascending=False)
    dsm_summary.insert(0, "Rank", range(1, len(dsm_summary) + 1))
    dsm_summary = dsm_summary.rename(columns={"Company Name": "Site"})
    
    dsm_summary = dsm_summary[["Rank", "Customer DSM", "Site", "WinBack_Lbs", "WinBack_Customers", 
                            "Conversion_Lbs", "Conversion_Customers", "Total_Opportunity_Lbs", "Total_Customers"]]
    
    # TERRITORY DETAIL
    territory_detail = pd.DataFrame()
    if "Customer Territory" in winback_opps.columns and "Customer Territory" in conversion_opps.columns:
        territory_winback = winback_opps.groupby(["Customer DSM", "Customer Territory"], dropna=False).agg(
            WinBack_Lbs=("Opportunity_Lbs", "sum")
        ).reset_index()
        
        territory_conversion = conversion_opps.groupby(["Customer DSM", "Customer Territory"], dropna=False).agg(
            Conversion_Lbs=("Opportunity_Lbs", "sum")
        ).reset_index()
        
        territory_detail = territory_winback.merge(territory_conversion, on=["Customer DSM", "Customer Territory"], how="outer").fillna(0)
        territory_detail["Total_Lbs"] = territory_detail["WinBack_Lbs"] + territory_detail["Conversion_Lbs"]
    
    # TOP TARGETS PER DSM
    winback_targets_list = []
    for dsm in dsm_summary["Customer DSM"].head(20):
        dsm_winbacks = winback_opps[winback_opps["Customer DSM"] == dsm].sort_values("Opportunity_Lbs", ascending=False).head(5)
        
        # Include all available columns including item details
        cols = [
            "Customer DSM", "Customer Name", "Company Customer Number", "Customer Territory",
            "Opportunity_Lbs", "Pounds_PY", "Pounds_CY", "YoY_Pct"
        ]
        if "Items_Lost" in dsm_winbacks.columns:
            cols.append("Items_Lost")
        
        winback_targets_list.append(dsm_winbacks[cols])

    conversion_targets_list = []
    for dsm in dsm_summary["Customer DSM"].head(20):
        dsm_conversions = conversion_opps[conversion_opps["Customer DSM"] == dsm].sort_values("Opportunity_Lbs", ascending=False).head(5)
        
        # Include all available columns including item details
        cols = [
            "Customer DSM", "Customer Name", "Company Customer Number", "Customer Territory",
            "Opportunity_Lbs", "Pounds_CY", "Conversion_Status"
        ]
        if "Items_To_Convert" in dsm_conversions.columns:
            cols.append("Items_To_Convert")
        
        conversion_targets_list.append(dsm_conversions[cols])
    
    winback_targets = pd.concat(winback_targets_list, ignore_index=True) if winback_targets_list else pd.DataFrame()
    conversion_targets = pd.concat(conversion_targets_list, ignore_index=True) if conversion_targets_list else pd.DataFrame()
    
    # EXPORT
    if not winback_opps.empty:
        try:
            winback_opps.to_csv("C:/Temp/WinBack_FULL_Export.csv", index=False)
            print(f"    ✓ Exported {len(winback_opps)} FULL win-back customers")
        except:
            pass
    
    print(f"    ✓ DSM Summary: {len(dsm_summary)} DSMs")
    
    return dsm_summary, territory_detail, winback_targets, conversion_targets

def generate_territory_routes(dsm_summary, winback_targets, conversion_targets, top_n_dsms=10):
    """Generate weekly route plans for top DSMs by territory. ONE visit per customer with item details."""
    
    if dsm_summary.empty:
        return {}
    
    print(f"\n🗺️ Building Territory Routes for Top {top_n_dsms} DSMs...")
    
    route_plans = {}
    top_dsms = dsm_summary.head(top_n_dsms)
    
    for idx, dsm_row in top_dsms.iterrows():
        dsm_name = dsm_row["Customer DSM"]
        rank = dsm_row["Rank"]
        
        print(f"  Processing Rank {rank}: {dsm_name}")
        
        # Get targets for this DSM
        dsm_winback = winback_targets[winback_targets["Customer DSM"] == dsm_name].copy() if not winback_targets.empty else pd.DataFrame()
        dsm_conversion = conversion_targets[conversion_targets["Customer DSM"] == dsm_name].copy() if not conversion_targets.empty else pd.DataFrame()
        
        # COMBINE customers (deduplicate by customer number)
        customer_aggregates = {}
        
        # Process win-back targets
        for _, row in dsm_winback.iterrows():
            cust_num = row["Company Customer Number"]
            territory = row.get("Customer Territory", "Unknown")
            key = (cust_num, territory)
            
            if key not in customer_aggregates:
                customer_aggregates[key] = {
                    "Customer Name": row["Customer Name"],
                    "Company Customer Number": cust_num,
                    "Customer Territory": territory,
                    "Target_Types": set(),
                    "Opportunity_Lbs": 0,
                    "Pounds_PY": row.get("Pounds_PY", 0),
                    "Pounds_CY": row.get("Pounds_CY", 0),
                    "Conversion_Status": "",
                    "Items_Lost": row.get("Items_Lost", ""),
                    "Items_To_Convert": "",
                }
            
            customer_aggregates[key]["Target_Types"].add("Win-Back")
            # For win-back, opportunity is the LOSS (PY - CY)
            winback_opp = abs(row.get("Pounds_PY", 0) - row.get("Pounds_CY", 0))
            customer_aggregates[key]["Opportunity_Lbs"] = max(
                customer_aggregates[key]["Opportunity_Lbs"],
                winback_opp
            )
        
        # Process conversion targets
        for _, row in dsm_conversion.iterrows():
            cust_num = row["Company Customer Number"]
            territory = row.get("Customer Territory", "Unknown")
            key = (cust_num, territory)
            
            if key not in customer_aggregates:
                customer_aggregates[key] = {
                    "Customer Name": row["Customer Name"],
                    "Company Customer Number": cust_num,
                    "Customer Territory": territory,
                    "Target_Types": set(),
                    "Opportunity_Lbs": row.get("Opportunity_Lbs", 0),
                    "Pounds_PY": 0,
                    "Pounds_CY": row.get("Pounds_CY", 0),
                    "Conversion_Status": row.get("Conversion_Status", ""),
                    "Items_Lost": "",
                    "Items_To_Convert": row.get("Items_To_Convert", ""),
                }
            else:
                # Customer is BOTH win-back and conversion
                customer_aggregates[key]["Target_Types"].add("Conversion")
                customer_aggregates[key]["Conversion_Status"] = row.get("Conversion_Status", "")
                customer_aggregates[key]["Items_To_Convert"] = row.get("Items_To_Convert", "")
                # Don't double-count pounds - use the larger opportunity
                customer_aggregates[key]["Opportunity_Lbs"] = max(
                    customer_aggregates[key]["Opportunity_Lbs"],
                    row.get("Opportunity_Lbs", 0)
                )
        
        # Convert to DataFrame with combined target types
        all_targets = []
        for key, data in customer_aggregates.items():
            target_type = " + ".join(sorted(data["Target_Types"]))
            all_targets.append({
                "Customer Name": data["Customer Name"],
                "Company Customer Number": data["Company Customer Number"],
                "Customer Territory": data["Customer Territory"],
                "Target_Type": target_type,
                "Opportunity_Lbs": data["Opportunity_Lbs"],
                "Pounds_PY": data["Pounds_PY"],
                "Pounds_CY": data["Pounds_CY"],
                "Conversion_Status": data["Conversion_Status"],
                "Items_Lost": data["Items_Lost"],
                "Items_To_Convert": data["Items_To_Convert"],
            })
        
        if not all_targets:
            continue
        
        targets_df = pd.DataFrame(all_targets)
        
        # Group by territory
        territories = targets_df["Customer Territory"].dropna().unique()
        
        for territory in territories:
            if str(territory).strip() == "" or str(territory).lower() == "unknown":
                continue
                
            territory_targets = targets_df[targets_df["Customer Territory"] == territory].copy()
            
            if territory_targets.empty:
                continue
            
            # Sort by opportunity descending (biggest opportunities first)
            territory_targets = territory_targets.sort_values("Opportunity_Lbs", ascending=False)
            
            # Calculate priority score
            territory_targets["Priority_Score"] = territory_targets["Opportunity_Lbs"]
            
            # Distribute across 5 days
            days = ["Monday", "Tuesday", "Wednesday", "Thursday", "Friday"]
            targets_per_day = max(1, len(territory_targets) // 5)
            
            territory_targets["Recommended_Day"] = ""
            territory_targets["Visit_Order"] = 0
            
            for day_idx, day in enumerate(days):
                start_idx = day_idx * targets_per_day
                end_idx = start_idx + targets_per_day if day_idx < 4 else len(territory_targets)
                
                for visit_order, idx in enumerate(range(start_idx, end_idx), start=1):
                    if idx < len(territory_targets):
                        territory_targets.iloc[idx, territory_targets.columns.get_loc("Recommended_Day")] = day
                        territory_targets.iloc[idx, territory_targets.columns.get_loc("Visit_Order")] = visit_order
            
            # Sort by day then visit order
            territory_targets = territory_targets.sort_values(["Recommended_Day", "Visit_Order"])
            
            # Create route key
            safe_territory = "".join(c for c in str(territory) if c.isalnum() or c in (' ', '-', '_')).strip()[:30]
            route_key = f"Rank{rank:02d}_{dsm_name.replace(' ', '')}_{safe_territory}"
            
            route_plans[route_key] = territory_targets
            
            total_opps = territory_targets["Opportunity_Lbs"].sum()
            print(f"    ✓ {territory}: {len(territory_targets)} customers, {total_opps:,.0f} lbs")
    
    print(f"\n  📋 Created {len(route_plans)} territory route plans")
    return route_plans

def write_dsm_scorecard_tab(xw, dsm_summary, territory_detail, winback_targets, conversion_targets):
    """
    Write the COMBINED DSM Opportunity Scorecard to Excel.
    Shows both win-back and conversion opportunities.
    """
    from openpyxl.styles import Font, PatternFill
    
    sheet_name = "13_DSM_Opportunity"
    
    # === SECTION 1: DSM SUMMARY (Combined) ===
    start_row = 0
    
    ws = xw.book.create_sheet(sheet_name)
    ws.cell(start_row + 1, 1).value = "DSM OPPORTUNITY SCORECARD (TRS ACCOUNTS ONLY)"
    ws.cell(start_row + 1, 1).font = Font(bold=True, size=14, color="0066CC")
    ws.merge_cells(f"A{start_row + 1}:H{start_row + 1}")
    
    ws.cell(start_row + 2, 1).value = "Win-Back = Active customers declining in category | Conversion = Buying but not aligned"
    ws.cell(start_row + 2, 1).font = Font(size=10, italic=True, color="666666")
    ws.merge_cells(f"A{start_row + 2}:H{start_row + 2}")
    
    dsm_summary.to_excel(xw, sheet_name=sheet_name, index=False, startrow=start_row + 3)
    
    _format_table_at(
        ws,
        header_row_0idx=start_row + 3,
        n_rows=dsm_summary.shape[0],
        number_headers={"WinBack_Lbs", "Conversion_Lbs", "Total_Opportunity_Lbs", "WinBack_Customers", "Conversion_Customers", "Total_Customers"},
        percent_headers=set()
    )
    
    # Highlight top 3 DSMs
    for row_idx in range(start_row + 5, min(start_row + 8, start_row + 4 + dsm_summary.shape[0] + 1)):
        for col_idx in range(1, 10):
            ws.cell(row_idx, col_idx).fill = PatternFill(
                start_color="FFF2CC", end_color="FFF2CC", fill_type="solid"
            )
    
    # === SECTION 2: TERRITORY DETAIL ===
    territory_start = start_row + dsm_summary.shape[0] + 6
    
    if not territory_detail.empty:
        ws.cell(territory_start, 1).value = "TERRITORY BREAKDOWN (Win-Back + Conversion by Territory)"
        ws.cell(territory_start, 1).font = Font(bold=True, size=12)
        
        territory_detail.to_excel(xw, sheet_name=sheet_name, index=False, startrow=territory_start + 1)
        
        _format_table_at(
            ws,
            header_row_0idx=territory_start + 1,
            n_rows=territory_detail.shape[0],
            number_headers={"WinBack_Lbs", "Conversion_Lbs", "Total_Lbs"},
            percent_headers=set()
        )
    
    # === SECTION 3: WIN-BACK TARGETS ===
    winback_start = territory_start + territory_detail.shape[0] + 4
    
    if not winback_targets.empty:
        ws.cell(winback_start, 1).value = "TOP 5 WIN-BACK TARGETS PER DSM (Declining but Active)"
        ws.cell(winback_start, 1).font = Font(bold=True, size=12, color="C00000")
        
        winback_targets.to_excel(xw, sheet_name=sheet_name, index=False, startrow=winback_start + 1)
        
        _format_table_at(
            ws,
            header_row_0idx=winback_start + 1,
            n_rows=winback_targets.shape[0],
            number_headers={"Opportunity_Lbs", "Pounds_PY", "Pounds_CY"},
            percent_headers={"YoY_Pct"}
        )
    
    # === SECTION 4: CONVERSION TARGETS ===
    conversion_start = winback_start + winback_targets.shape[0] + 4
    
    if not conversion_targets.empty:
        ws.cell(conversion_start, 1).value = "TOP 5 CONVERSION TARGETS PER DSM (Buying but Not Aligned)"
        ws.cell(conversion_start, 1).font = Font(bold=True, size=12, color="E67E22")
        
        conversion_targets.to_excel(xw, sheet_name=sheet_name, index=False, startrow=conversion_start + 1)
        
        _format_table_at(
            ws,
            header_row_0idx=conversion_start + 1,
            n_rows=conversion_targets.shape[0],
            number_headers={"Opportunity_Lbs", "Pounds_CY"},
            percent_headers=set()
        )
    
    print("  ✓ Combined DSM Opportunity Scorecard created (Win-Back + Conversion)")

def create_territory_routes_workbook(route_plans_dict, outdir, fiscal_year, current_week):
    """
    Create a separate Excel workbook for territory routes.
    Similar to Sales_Leads.csv but organized by territory with route schedules.
    
    Args:
        route_plans_dict: Dictionary of {route_key: route_dataframe}
        outdir: Output directory
        fiscal_year: Fiscal year for filename
        current_week: Current week for filename
    
    Returns:
        Path to created workbook
    """
    from openpyxl.styles import Font, PatternFill, Alignment
    from datetime import datetime
    
    if not route_plans_dict:
        print("    ℹ No territory routes to create")
        return None
    
    print(f"\n📗 Creating Territory Routes Workbook...")
    
    # Create filename
    stamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    workbook_name = f"Territory_Routes_FY{fiscal_year}_wk{current_week}_{stamp}.xlsx"
    workbook_path = os.path.join(outdir, workbook_name)
    
    # Day colors
    day_colors = {
        "Monday": "E8F4F8",
        "Tuesday": "FFF4E6",
        "Wednesday": "F0F8E8",
        "Thursday": "FCE8F3",
        "Friday": "F3E8FC"
    }
    
    with pd.ExcelWriter(workbook_path, engine="openpyxl") as writer:
        
        # Create INDEX tab first
        index_data = []
        for route_key, route_df in route_plans_dict.items():
            parts = route_key.split("_", 2)
            rank = parts[0].replace("Rank", "") if len(parts) > 0 else "?"
            dsm = parts[1] if len(parts) > 1 else "Unknown"
            territory = parts[2] if len(parts) > 2 else "Unknown"
            
            total_visits = len(route_df)
            total_opps = route_df["Opportunity_Lbs"].sum()
            winback = (route_df["Target_Type"] == "Win-Back").sum() if "Target_Type" in route_df.columns else 0
            conversion = (route_df["Target_Type"] == "Conversion").sum() if "Target_Type" in route_df.columns else 0
            
            index_data.append({
                "Rank": int(rank),
                "DSM": dsm,
                "Territory": territory,
                "Total_Visits": total_visits,
                "Win_Back": winback,
                "Conversion": conversion,
                "Total_Opportunity_Lbs": total_opps,
                "Tab_Name": route_key[:31]
            })
        
        index_df = pd.DataFrame(index_data).sort_values("Rank")
        index_df.to_excel(writer, sheet_name="Index", index=False)
        
        # Format Index tab
        ws_index = writer.book["Index"]
        ws_index.insert_rows(1, 3)
        ws_index.cell(1, 1, value="TERRITORY ROUTE PLANS INDEX").font = Font(bold=True, size=16, color="0066CC")
        ws_index.merge_cells("A1:H1")
        ws_index.cell(2, 1, value=f"Generated: {datetime.now():%B %d, %Y at %I:%M %p} | FY{fiscal_year} Week {current_week}").font = Font(size=10, italic=True, color="666666")
        ws_index.merge_cells("A2:H2")
        ws_index.cell(3, 1, value=f"Total Routes: {len(route_plans_dict)} | Click tab names below to view schedules").font = Font(size=11, color="000000")
        ws_index.merge_cells("A3:H3")
        
        # Format index numbers
        for col_name in ["Rank", "Total_Visits", "Win_Back", "Conversion", "Total_Opportunity_Lbs"]:
            col_idx = None
            for cell in ws_index[4]:
                if cell.value == col_name:
                    col_idx = cell.col_idx
                    break
            if col_idx:
                for row_idx in range(5, 5 + len(index_df)):
                    cell = ws_index.cell(row_idx, col_idx)
                    if cell.value is not None:
                        try:
                            cell.value = float(str(cell.value).replace(',', ''))
                            cell.number_format = '#,##0'
                        except:
                            pass
        
        # Auto-size Index columns
        for col_idx in range(1, ws_index.max_column + 1):
            max_length = 0
            column_letter = get_column_letter(col_idx)
            
            for row in ws_index.iter_rows(min_col=col_idx, max_col=col_idx):
                cell = row[0]
                try:
                    if cell.value:
                        max_length = max(max_length, len(str(cell.value)))
                except:
                    pass
            
            ws_index.column_dimensions[column_letter].width = min(max_length + 2, 50)
            
            for row in ws_index.iter_rows(min_col=col_idx, max_col=col_idx):
                cell = row[0]
                try:
                    if cell.value:
                        max_length = max(max_length, len(str(cell.value)))
                except:
                    pass
            
            ws_index.column_dimensions[column_letter].width = min(max_length + 2, 50)
        
        # Create individual route tabs
        for idx, (route_key, route_df) in enumerate(route_plans_dict.items(), start=1):
            parts = route_key.split("_", 2)
            rank_part = parts[0] if len(parts) > 0 else "00"
            dsm_part = parts[1] if len(parts) > 1 else "Unknown"
            territory_part = parts[2] if len(parts) > 2 else "Territory"
            
            # Tab name (31 char limit)
            sheet_name = route_key[:31]
            
            # Write route data
            route_df.to_excel(writer, sheet_name=sheet_name, index=False, startrow=4)
            
            ws = writer.book[sheet_name]
            
            # Header
            ws.cell(1, 1, value=f"Weekly Route: {territory_part}").font = Font(bold=True, size=14, color="0066CC")
            ws.merge_cells("A1:H1")
            
            # Summary
            total_opps = route_df["Opportunity_Lbs"].sum()
            total_visits = len(route_df)
            winback_count = (route_df["Target_Type"] == "Win-Back").sum() if "Target_Type" in route_df.columns else 0
            conversion_count = (route_df["Target_Type"] == "Conversion").sum() if "Target_Type" in route_df.columns else 0
            
            ws.cell(2, 1, value=f"DSM: {dsm_part} | Territory Manager: {territory_part}").font = Font(size=11, bold=True)
            ws.merge_cells("A2:H2")
            
            ws.cell(3, 1, value=f"Total: {total_visits} visits | {total_opps:,.0f} lbs opportunity | Win-Back: {winback_count} | Conversion: {conversion_count}").font = Font(size=10, color="666666")
            ws.merge_cells("A3:H3")
            
            # Format numbers
            header_row = 5
            col_map = {}
            for cell in ws[header_row]:
                if cell.value:
                    col_map[str(cell.value).strip()] = cell.col_idx
            
            for col_name in ["Opportunity_Lbs", "Priority_Score", "Pounds_PY", "Pounds_CY", "Visit_Order"]:
                if col_name in col_map:
                    col_idx = col_map[col_name]
                    for row_idx in range(header_row + 1, header_row + 1 + len(route_df)):
                        cell = ws.cell(row_idx, col_idx)
                        if cell.value is not None:
                            try:
                                cell.value = float(str(cell.value).replace(',', ''))
                                cell.number_format = '#,##0'
                            except:
                                pass
            
            # Color by day
            if "Recommended_Day" in col_map:
                day_col = col_map["Recommended_Day"]
                for row_idx in range(header_row + 1, header_row + 1 + len(route_df)):
                    day_value = ws.cell(row_idx, day_col).value
                    if day_value in day_colors:
                        for col_idx in range(1, len(route_df.columns) + 1):
                            ws.cell(row_idx, col_idx).fill = PatternFill(
                                start_color=day_colors[day_value],
                                end_color=day_colors[day_value],
                                fill_type="solid"
                            )
            
            # Auto-size
            for col_idx in range(1, ws.max_column + 1):
                max_length = 0
                column_letter = get_column_letter(col_idx)
                
                for row in ws.iter_rows(min_col=col_idx, max_col=col_idx):
                    cell = row[0]
                    try:
                        if cell.value:  # ← SIMPLER - just check if there's a value
                            max_length = max(max_length, len(str(cell.value)))
                    except:
                        pass
                
                adjusted_width = min(max_length + 2, 50)
                ws.column_dimensions[column_letter].width = adjusted_width
            
            print(f"    ✓ Tab {idx}/{len(route_plans_dict)}: {territory_part}")
    
    print(f"  ✅ Territory Routes Workbook: {workbook_name}")
    return workbook_path


def write_site_lead_summary_tab(xw, leads_df):
    """
    Create dedicated tab showing sites ranked by lead volume.
    Simple summary for quick reference.
    """
    from openpyxl.styles import Font, PatternFill
    
    if leads_df.empty:
        return
    
    sheet_name = "12_Sites_by_Leads"
    
    # Aggregate by site
    site_summary = leads_df.groupby("Company Name", dropna=False).agg(
        Total_Leads=("Company Customer Number", "nunique"),
        Total_Opportunity_Lbs=("Pounds_CY", "sum"),
        Avg_Lbs_per_Lead=("Pounds_CY", "mean")
    ).reset_index().sort_values("Total_Opportunity_Lbs", ascending=False)
    
    site_summary.insert(0, "Rank", range(1, len(site_summary) + 1))
    
    # Add account type breakdown if available
    if "Customer Account Type Code" in leads_df.columns:
        acct_breakdown = leads_df.groupby(["Company Name", "Customer Account Type Code"]).size().unstack(fill_value=0)
        acct_breakdown.columns = [f"{col}_Leads" for col in acct_breakdown.columns]
        site_summary = site_summary.merge(acct_breakdown, left_on="Company Name", right_index=True, how="left")
    
    ws = xw.book.create_sheet(sheet_name)
    
    ws.cell(1, 1).value = "SITES RANKED BY SALES LEAD VOLUME"
    ws.cell(1, 1).font = Font(bold=True, size=14, color="0066CC")
    ws.merge_cells(f"A1:F1")
    
    ws.cell(2, 1).value = "Sites with the most actionable sales leads (conversion + win-back opportunities)"
    ws.cell(2, 1).font = Font(size=10, italic=True, color="666666")
    ws.merge_cells(f"A2:F2")
    
    site_summary.to_excel(xw, sheet_name=sheet_name, index=False, startrow=3)
    
    number_cols = {"Total_Leads", "Total_Opportunity_Lbs", "Avg_Lbs_per_Lead"}
    if "TRS_Leads" in site_summary.columns:
        number_cols.add("TRS_Leads")
    if "LCC_Leads" in site_summary.columns:
        number_cols.add("LCC_Leads")
    if "CMU_Leads" in site_summary.columns:
        number_cols.add("CMU_Leads")
    
    _format_table_at(
        ws,
        header_row_0idx=3,
        n_rows=site_summary.shape[0],
        number_headers=number_cols,
        percent_headers=set()
    )
    
    # Highlight top 5 sites
    for row_idx in range(5, min(10, 4 + site_summary.shape[0] + 1)):
        for col_idx in range(1, 8):
            ws.cell(row_idx, col_idx).fill = PatternFill(
                start_color="D9EAD3", end_color="D9EAD3", fill_type="solid"
            )
    
    print("  ✓ Sites by Lead Volume tab created")

def _autosize_columns(ws, min_width=10, max_width=50):
    """Auto-fit column widths based on content."""
    from openpyxl.utils import get_column_letter
    
    # Auto-size columns (safe version)
    for col_idx in range(1, ws.max_column + 1):
        max_length = 0
        column_letter = get_column_letter(col_idx)
        
        for row in ws.iter_rows(min_col=col_idx, max_col=col_idx):
            cell = row[0]
            try:
                if cell.value:
                    max_length = max(max_length, len(str(cell.value)))
            except:
                pass
        
        adjusted_width = min(max_length + 2, 50)
        ws.column_dimensions[column_letter].width = adjusted_width

def _autosize_rows(ws):
    """Auto-adjust row heights to fit wrapped text."""
    from openpyxl.utils import get_column_letter
    
    for row_idx in range(1, ws.max_row + 1):
        max_height = 15  # Minimum height in points
        
        for col_idx in range(1, ws.max_column + 1):
            cell = ws.cell(row_idx, col_idx)
            
            if not cell.value:
                continue
            
            # Only calculate for wrapped text
            if cell.alignment and cell.alignment.wrap_text:
                # Get column width in characters
                col_letter = get_column_letter(col_idx)
                col_width = ws.column_dimensions[col_letter].width or 10
                
                # Convert to approximate character count
                chars_per_line = int(col_width * 1.0)  # Rough estimate
                
                # Count lines needed
                text = str(cell.value)
                text_length = len(text)
                
                # Account for explicit line breaks
                explicit_lines = text.count('\n') + 1
                
                # Estimate wrapped lines
                wrapped_lines = max(1, text_length / chars_per_line) if chars_per_line > 0 else 1
                
                # Total lines needed
                total_lines = max(explicit_lines, wrapped_lines)
                
                # Calculate height (approximately 15 points per line)
                estimated_height = total_lines * 15
                
                # Cap at reasonable maximum
                estimated_height = min(estimated_height, 150)
                
                max_height = max(max_height, estimated_height)
        
        # Apply the calculated height
        ws.row_dimensions[row_idx].height = max_height
        
def write_grouped_sales_leads(xw, leads_df, sheet_name="12_Sales_Leads"):
    """
    Write hierarchical sales leads: DSM Summary -> Customer Detail
    (Site summary removed per user request)
    """
    if leads_df.empty:
        pd.DataFrame({"Note": ["No sales leads found"]}).to_excel(xw, sheet_name=sheet_name, index=False)
        return
    
    # Create the sheet
    ws = xw.book.create_sheet(sheet_name)
    start_row = 0
    
    # === SECTION 1: DSM SUMMARY ===
    if "Customer DSM" in leads_df.columns:
        dsm_summary = leads_df.groupby(["Company Name", "Customer DSM"], dropna=False).agg(
            Num_Leads=("Company Customer Number", "nunique"),
            Opportunity_Lbs=("Pounds_CY", "sum")
        ).reset_index().sort_values(["Company Name", "Opportunity_Lbs"], ascending=[True, False])
        
        ws.cell(start_row + 1, 1).value = "DSM SALES OPPORTUNITIES (Grouped by Site)"
        ws.cell(start_row + 1, 1).font = Font(bold=True, size=14, color="0066CC")
        ws.merge_cells(f"A{start_row + 1}:D{start_row + 1}")
        
        dsm_summary.to_excel(xw, sheet_name=sheet_name, index=False, startrow=start_row + 2)
        
        _format_table_at(
            ws,
            header_row_0idx=start_row + 2,
            n_rows=dsm_summary.shape[0],
            number_headers={"Num_Leads", "Opportunity_Lbs"},
            percent_headers=set()
        )
        
        detail_start = start_row + dsm_summary.shape[0] + 5
    else:
        detail_start = start_row + 2
    
    # === SECTION 2: DETAILED LEADS (Grouped by Site, then DSM) ===
    ws.cell(detail_start, 1).value = "DETAILED CUSTOMER LEADS (Grouped by Site & DSM)"
    ws.cell(detail_start, 1).font = Font(bold=True, size=12)
    
    # Sort leads for logical grouping
    leads_sorted = leads_df.sort_values(["Company Name", "Customer DSM", "Pounds_CY"], ascending=[True, True, False])
    
    leads_sorted.to_excel(xw, sheet_name=sheet_name, index=False, startrow=detail_start + 1)
    
    _format_table_at(
        ws,
        header_row_0idx=detail_start + 1,
        n_rows=leads_sorted.shape[0],
        number_headers={"Pounds_CY"},
        percent_headers=set()
    )
    
    print("  ✓ Grouped Sales Leads tab created (Site -> DSM -> Customer hierarchy)")
        
# ---------------------------- Write workbook (extended) ----------------------------
def write_excel(
    excel_path,
    summary_kpi, lag_sites, cust_losses,
    status, leads, vendor_index,
    company_yoy=None,
    all_weekly=None, trs_weekly=None,
    cmu_weekly=None, lcc_weekly=None,
    forecast_method="linear",
    status_raw=None,
    alignment_df=None,
    filtered_company_name=None,
    min_ytd=260,
    source_path=None,
    dsm_summary=None,              
    territory_detail=None,         
    winback_targets=None,          
    conversion_targets=None       
):
    # Debug output
    print(f"\n  Starting Excel write to: {excel_path}")
    print(f"  File parent directory exists: {os.path.exists(os.path.dirname(excel_path))}")
    
    # Get fiscal period info once - USE status_raw NOT raw_df
    min_week, max_week, fiscal_year = get_fiscal_period_info(status_raw)
    weeks_covered = max_week - min_week + 1
    
    try:
        forecast_desc = "Linear Trend Forecast" if (forecast_method or "linear").lower() == "linear" \
                        else "Run-Rate Forecast (Last 4 Weeks Avg)"

        with pd.ExcelWriter(excel_path, engine="openpyxl") as xw:
            # Calculate weeks covered
            if status_raw is not None and 'Fiscal Week Number' in status_raw.columns:
                weeks_series = pd.to_numeric(status_raw['Fiscal Week Number'], errors='coerce').dropna()
                if len(weeks_series) > 0:
                    min_week = int(weeks_series.min())
                    max_week = int(weeks_series.max())
                    weeks_covered = max_week - min_week + 1
                else:
                    min_week = 1
                    max_week = 52
                    weeks_covered = 52
            else:
                min_week = 1
                max_week = 52
                weeks_covered = 52
            
            # CREATE GUIDE TAB
            create_guide_tab(xw, status, weeks_covered, filtered_company_name, min_week, max_week)
            
            # SUMMARY TAB
            summary_kpi.to_excel(xw, sheet_name="Summary", index=False, startrow=0)
            

            explain = pd.DataFrame({"What this shows":[
                "Conversion KPIs: share of CY pounds aligned to SUPC+SUVC.",
                "Lagging Sites: Companies with the largest negative Δ Pounds YoY and their share of total category losses.",
                "Top Customer Losses: Customers with the biggest YoY declines to target for recovery.",
                "Recommendations: Auto-generated next steps based on losses and conversion gaps."
            ]})
            explain_start = summary_kpi.shape[0] + 2
            explain.to_excel(xw, sheet_name="Summary", index=False, startrow=explain_start)
            
            # Overall YoY metrics
            overall_start = explain_start + explain.shape[0] + 3
            ws_summary = xw.book["Summary"]
            ws_summary.cell(overall_start, 1, value="All Account Type Volume YoY").font = Font(bold=True, size=12)

            overall_yoy = _overall_yoy(status)
            overall_yoy.to_excel(xw, sheet_name="Summary", index=False, startrow=overall_start+1)
            _format_table_at(
                ws_summary,
                header_row_0idx=overall_start+1,
                n_rows=overall_yoy.shape[0],
                number_headers={"Pounds CY", "Pounds PY", "Delta Pounds YoY"},
                percent_headers={"YoY %"}
            )

            # Sysco Brand split
            brand_start = overall_start + overall_yoy.shape[0] + 4
            ws_summary.cell(brand_start, 1, value="Sysco Brand YoY, All Account Types").font = Font(bold=True, size=12)

            if status_raw is not None and "Sysco Brand Indicator" in status_raw.columns:
                # Aggregate raw data by brand
                brand_agg = status_raw.groupby("Sysco Brand Indicator", dropna=False).agg({
                    "Pounds CY": "sum",
                    "Pounds PY": "sum"
                }).reset_index()
                brand_agg.columns = ["Sysco Brand Indicator", "Pounds_CY", "Pounds_PY"]
                brand_agg["Delta_YoY_Lbs"] = brand_agg["Pounds_CY"] - brand_agg["Pounds_PY"]
                brand_agg["YoY_Pct"] = np.where(
                    brand_agg["Pounds_PY"] > 0,
                    brand_agg["Delta_YoY_Lbs"] / brand_agg["Pounds_PY"],
                    np.nan
                )
                brand_split = brand_agg
            else:
                brand_split = pd.DataFrame(columns=["Sysco Brand Indicator","Pounds_CY","Pounds_PY","Delta_YoY_Lbs","YoY_Pct"])

            brand_split.to_excel(xw, sheet_name="Summary", index=False, startrow=brand_start+1)
            _format_table_at(
                ws_summary,
                header_row_0idx=brand_start+1,
                n_rows=brand_split.shape[0],
                number_headers={"Pounds_CY", "Pounds_PY", "Delta_YoY_Lbs"},
                percent_headers={"YoY_Pct"}
            )

            # (re)add lag tables so Summary has CY/PY/Delta columns visible on that sheet
            lag_out = lag_sites.rename(columns={"d":"Delta_YoY_Lbs"})
            lag_start = brand_start + brand_split.shape[0] + 3
            lag_out.to_excel(xw, sheet_name="Summary", index=False, startrow=lag_start)

            cust_out = cust_losses.rename(columns={"d":"Delta_YoY_Lbs"})
            cust_start = lag_start + lag_out.shape[0] + 3
            cust_out.to_excel(xw, sheet_name="Summary", index=False, startrow=cust_start)

            # Format each table on Summary individually
            _format_table_at(
                xw.book["Summary"],
                header_row_0idx=0,
                n_rows=summary_kpi.shape[0],
                number_headers={"Value", "Customers CY", "Customers PY"},  # ADDED
                percent_headers={"% Any Item aligned", "% Any Vendor aligned", "% Item+Vendor aligned", "% Item-only", "% Vendor-only", "% Neither", "Customer Retention %"}  # ADDED
            )

            _format_table_at(
                xw.book["Summary"],
                header_row_0idx=lag_start,  # lag_sites table
                n_rows=lag_out.shape[0],
                number_headers={"CY", "PY", "Delta_YoY_Lbs"},
                percent_headers={"Loss_%_of_total"}
            )

            _format_table_at(
                xw.book["Summary"],
                header_row_0idx=cust_start,  # customer losses table
                n_rows=cust_out.shape[0],
                number_headers={"CY", "PY", "Delta_YoY_Lbs"},
                percent_headers=set()
            )

            # ===== PREP for ACCOUNT TABS =====
            base_cols = [
                "Pounds_CY","Pounds_PY","Delta_YoY_Lbs","YoY_Pct",
                "Sysco Brand Indicator","Company Name",
                "Item Number","Item Description",
                "Customer Account Type Code"
            ]
            if status_raw is not None:
                for c in base_cols:
                    if c not in status.columns and c in status_raw.columns:
                        status[c] = status_raw[c]

            def slice_acct(df, codes=None):
                if "Customer Account Type Code" not in df.columns or not codes:
                    return df
                keep = set(x.strip() for x in codes)
                return df[df["Customer Account Type Code"].astype(str).isin(keep)].copy()

            # ===== ACCOUNT TABS with filtered_company_name =====
            _write_account_tab(xw, "01_All_Accounts", status.copy(), filtered_company=filtered_company_name)
            _write_account_tab(xw, "02_TRS", slice_acct(status, {"TRS"}), filtered_company=filtered_company_name)
            _write_account_tab(xw, "03_LCC", slice_acct(status, {"LCC"}), filtered_company=filtered_company_name)
            _write_account_tab(xw, "04_CMU", slice_acct(status, {"CMU"}), filtered_company=filtered_company_name)

            # 05b / 07b — chart-only sheets
            print(f"  Checking chart data:")
            print(f"    all_weekly: exists={all_weekly is not None}, rows={len(all_weekly) if all_weekly is not None else 0}")
            print(f"    trs_weekly: exists={trs_weekly is not None}, rows={len(trs_weekly) if trs_weekly is not None else 0}")

            # ===== CHART TABS with filtered_company_name =====
            if all_weekly is not None and not all_weekly.empty:
                _write_weekly_chart_only_tab(xw, "05_All_Weekly_Chart", all_weekly, 
                                            forecast_method_desc=forecast_desc, 
                                            filtered_company=filtered_company_name)

            if trs_weekly is not None and not trs_weekly.empty:
                _write_weekly_chart_only_tab(xw, "07_TRS_Weekly_Chart", trs_weekly, 
                                            forecast_method_desc=forecast_desc, 
                                            filtered_company=filtered_company_name)

            if cmu_weekly is not None and not cmu_weekly.empty:
                _write_weekly_chart_only_tab(xw, "08_CMU_Weekly_Chart", cmu_weekly, 
                                            forecast_method_desc=forecast_desc, 
                                            filtered_company=filtered_company_name)

            if lcc_weekly is not None and not lcc_weekly.empty:
                _write_weekly_chart_only_tab(xw, "09_LCC_Weekly_Chart", lcc_weekly, 
                                            forecast_method_desc=forecast_desc, 
                                            filtered_company=filtered_company_name)
            

            # ===== 10 — Vendor/Item alignment =====
            align_view = leads.copy()
            align_view.to_excel(xw, sheet_name="10_VendorItem_Alignment", index=False)
            _try_format(xw, "08_VendorItem_Alignment",
                        number_headers={"Pounds_CY"},
                        percent_headers=set())
            
            # ===== 11 — Award vs Sales =====
            if alignment_df is not None and "Award Volume Annualized" in alignment_df.columns:
                # Get award volumes from alignment
                alignment_awards = alignment_df[["SUVC", "Supplier Name", "SUPC", "Award Volume Annualized"]].copy()
                alignment_awards["Award Volume Annualized"] = pd.to_numeric(
                    alignment_awards["Award Volume Annualized"], errors="coerce"
                ).fillna(0)
                
                # Sum by vendor
                vendor_awards = alignment_awards.groupby(["SUVC", "Supplier Name"], dropna=False).agg(
                    Award_Volume_Annualized=("Award Volume Annualized", "sum")
                ).reset_index()
                
                # Get sales from status_raw
                by_vendor = status_raw.groupby(["SUVC","Supplier Name"], dropna=False).agg(
                    CY_Lbs=("Pounds CY","sum"),
                    PY_Lbs=("Pounds PY","sum")
                ).reset_index()
                
                # Merge and calculate
                by_vendor = by_vendor.merge(vendor_awards, on=["SUVC", "Supplier Name"], how="left")
                by_vendor["Award_Volume_Annualized"] = by_vendor["Award_Volume_Annualized"].fillna(0)
                by_vendor["Expected_At_This_Point"] = by_vendor["Award_Volume_Annualized"] * (weeks_covered / 52.0)
                by_vendor["% of Expected"] = np.where(
                    by_vendor["Expected_At_This_Point"] > 0,
                    by_vendor["CY_Lbs"] / by_vendor["Expected_At_This_Point"], 
                    np.nan
                )
                by_vendor["% of Annual Award"] = np.where(
                    by_vendor["Award_Volume_Annualized"] > 0,
                    by_vendor["CY_Lbs"] / by_vendor["Award_Volume_Annualized"], 
                    np.nan
                )
                
                # Reorder columns
                by_vendor = by_vendor[[
                    "SUVC", "Supplier Name", 
                    "CY_Lbs", "PY_Lbs",
                    "Award_Volume_Annualized", "Expected_At_This_Point",
                    "% of Expected", "% of Annual Award"
                ]]
                
                # Write to Excel
                by_vendor.to_excel(xw, sheet_name="11_Award_vs_Sales", index=False)
                ws_award = xw.book["11_Award_vs_Sales"]
                ws_award.insert_rows(1, 3)
                
                ws_award.cell(1, 1, value=f"Report Coverage: Weeks {min_week}-{max_week} ({weeks_covered} of 52 weeks = {weeks_covered/52:.1%} of year)").font = Font(bold=True, size=12, color="0066CC")
                ws_award.cell(2, 1, value=f"✓ Use '% of Expected' column to evaluate performance (compares to pro-rated goal)").font = Font(size=10, color="008000")
                ws_award.cell(3, 1, value=f"ℹ '% of Annual Award' shows full-year pace (only valid for 52-week reports)").font = Font(size=10, color="666666")
                
                # Format numbers and percentages
                header_row = 5
                col_map = {}
                for cell in ws_award[header_row]:
                    if cell.value:
                        col_map[str(cell.value).strip()] = cell.col_idx
                
                # Number columns
                for col_name in ["CY_Lbs", "PY_Lbs", "Award_Volume_Annualized", "Expected_At_This_Point"]:
                    if col_name in col_map:
                        col_idx = col_map[col_name]
                        for row_idx in range(header_row + 1, header_row + 1 + by_vendor.shape[0]):
                            cell = ws_award.cell(row_idx, col_idx)
                            try:
                                if cell.value is not None:
                                    cell.value = float(str(cell.value).replace(',', ''))
                                    cell.number_format = '#,##0.0'
                            except:
                                pass
                
                # Percent columns
                for col_name in ["% of Expected", "% of Annual Award"]:
                    if col_name in col_map:
                        col_idx = col_map[col_name]
                        for row_idx in range(header_row + 1, header_row + 1 + by_vendor.shape[0]):
                            cell = ws_award.cell(row_idx, col_idx)
                            try:
                                if cell.value is not None:
                                    cell.value = float(cell.value)
                                    cell.number_format = '0.0%'
                            except:
                                pass
                
                # Conditional formatting
                from openpyxl.styles import PatternFill
                if "% of Expected" in col_map:
                    col_idx = col_map["% of Expected"]
                    for row in range(header_row + 1, header_row + 1 + by_vendor.shape[0]):
                        cell = ws_award.cell(row, col_idx)
                        try:
                            val = float(cell.value) if cell.value else 0
                            if val >= 1.0:
                                cell.fill = PatternFill(start_color="C6EFCE", end_color="C6EFCE", fill_type="solid")
                            elif val < 0.8:
                                cell.fill = PatternFill(start_color="FFC7CE", end_color="FFC7CE", fill_type="solid")
                        except:
                            pass

            else:
                pd.DataFrame({"Note":[
                    "Award Volume Annualized not found. Confirm Alignment.csv carries it AND the join keys match."
                ]}).to_excel(xw, sheet_name="11_Award_vs_Sales", index=False)

            # ===== 12 — SITES BY LEAD VOLUME =====
            write_site_lead_summary_tab(xw, leads)

            # ===== 13 — COMBINED DSM OPPORTUNITY SCORECARD =====
            print("\n📊 Building Combined DSM Opportunity Scorecard...")
            if dsm_summary is not None and not dsm_summary.empty:
                write_dsm_scorecard_tab(xw, dsm_summary, territory_detail, winback_targets, conversion_targets)
                
                # Export win-back details for validation
                if winback_targets is not None and not winback_targets.empty:
                    export_dir = os.path.dirname(excel_path)
                    export_path = os.path.join(export_dir, "WinBack_Detail_Export.csv")
                    winback_targets.to_csv(export_path, index=False)
                
                print(f"  ✓ DSM Opportunity Scorecard: {len(dsm_summary)} DSMs")
            else:
                print("  ℹ No DSM opportunities to write")

            vendor_index.to_excel(xw, sheet_name="14_Vendor_Leads_Index", index=False)

            # ===== 14 — NPD CUISINE TYPE ANALYSIS (TRS ONLY) =====
            print("  Creating NPD Cuisine Type analysis tab...")

            # Use status_raw which has NPD Cuisine Type column
            if status_raw is not None:
                trs_data = status_raw[status_raw["Customer Account Type Code"] == "TRS"].copy()
                
                if not trs_data.empty:
                    cuisine_analysis = _cuisine_customer_analysis(trs_data)
                    cuisine_analysis.to_excel(xw, sheet_name="14_NPD_Cuisine_TRS", index=False)
                    
                    # Format
                    ws_cuisine = xw.book["14_NPD_Cuisine_TRS"]
                    ws_cuisine.insert_rows(1, 2)
                    ws_cuisine.cell(1, 1, value="NPD Cuisine Type Customer Analysis - TRS Accounts Only").font = Font(bold=True, size=14, color="0066CC")
                    ws_cuisine.cell(2, 1, value="Which cuisine types are losing the most customers?").font = Font(size=10, italic=True, color="666666")
                    
                    if "Note" not in cuisine_analysis.columns:
                        # Build column map
                        header_row = 3
                        col_map = {}
                        for cell in ws_cuisine[header_row]:
                            if cell.value:
                                col_map[str(cell.value).strip()] = cell.col_idx
                        
                        # Format number columns
                        number_cols = ["Customers_CY", "Customers_PY", "Delta_Customers"]
                        for col_name in number_cols:
                            if col_name in col_map:
                                col_idx = col_map[col_name]
                                for row_idx in range(header_row + 1, header_row + 1 + cuisine_analysis.shape[0]):
                                    cell = ws_cuisine.cell(row_idx, col_idx)
                                    if cell.value is not None:
                                        try:
                                            cell.value = int(float(str(cell.value).replace(',', '')))
                                            cell.number_format = '#,##0'
                                        except:
                                            pass
                        
                        # Format percent column
                        if "Pct_Change" in col_map:
                            col_idx = col_map["Pct_Change"]
                            for row_idx in range(header_row + 1, header_row + 1 + cuisine_analysis.shape[0]):
                                cell = ws_cuisine.cell(row_idx, col_idx)
                                if cell.value is not None:
                                    try:
                                        cell.value = float(cell.value)
                                        cell.number_format = '0.0%'
                                    except:
                                        pass
                
            vendor_index.to_excel(xw, sheet_name="15_Vendor_Leads_Index", index=False)

            # Auto-size all sheets at the end
            print("  Formatting worksheets...")
            for sheet_name in xw.book.sheetnames:
                ws = xw.book[sheet_name]
                _autosize_columns(ws)
                _autosize_rows(ws)  # Add this line
            print("  ✓ Auto-sized columns and rows")

    except Exception as e:
        print(f"\n  ❌ EXCEL WRITE FAILED: {e}")
        import traceback
        traceback.print_exc()
        raise

def export_all_chart_images(excel_path):
    """
    Export all chart images from Excel workbook to Chart_Images folder.
    Call this AFTER Excel file is closed.
    """
    import platform
    if platform.system() != "Windows":
        print("    ⚠ Chart export only supported on Windows")
        return False
    
    try:
        import win32com.client
        
        chart_dir = os.path.join(os.path.dirname(excel_path), "Chart_Images")
        os.makedirs(chart_dir, exist_ok=True)
        
        print(f"    📸 Exporting chart images...")
        
        # Get or create Excel instance
        try:
            excel_app = win32com.client.GetActiveObject("Excel.Application")
        except:
            excel_app = win32com.client.Dispatch("Excel.Application")
        
        excel_app.DisplayAlerts = False
        excel_app.Visible = False
        
        # Open the workbook
        wb = excel_app.Workbooks.Open(os.path.abspath(excel_path))
        
        # Chart sheet names
        chart_sheets = [
            "05_All_Weekly_Chart",
            "07_TRS_Weekly_Chart", 
            "08_CMU_Weekly_Chart",
            "09_LCC_Weekly_Chart"
        ]
        
        exported = 0
        for sheet_name in chart_sheets:
            try:
                ws = wb.Worksheets(sheet_name)
                if ws.ChartObjects().Count > 0:
                    chart_path = os.path.join(chart_dir, f"{sheet_name}.png")
                    chart_obj = ws.ChartObjects(1)
                    
                    # PRIMARY METHOD: Clipboard (most reliable for all chart types)
                    try:
                        from PIL import ImageGrab
                        
                        # Copy chart to clipboard
                        chart_obj.CopyPicture(1, 2)  # xlScreen, xlBitmap
                        
                        # Small delay for clipboard stability
                        import time
                        time.sleep(0.2)
                        
                        # Grab and save
                        img = ImageGrab.grabclipboard()
                        if img:
                            img.save(chart_path, 'PNG')
                            print(f"      ✓ {sheet_name}.png")
                            exported += 1
                        else:
                            raise Exception("Clipboard empty")
                            
                    except ImportError:
                        # Pillow not installed - use fallback
                        print(f"      ⚠ Pillow not installed, using direct export (less reliable)")
                        chart_obj.Chart.Export(chart_path, "PNG")
                        print(f"      ✓ {sheet_name}.png (fallback)")
                        exported += 1
                        
                    except Exception as e:
                        # Clipboard failed - try direct export
                        print(f"      ⚠ Clipboard method failed, trying direct export: {e}")
                        chart_obj.Chart.Export(chart_path, "PNG")
                        print(f"      ✓ {sheet_name}.png (fallback)")
                        exported += 1
                        
            except Exception as e:
                print(f"      ❌ All methods failed for {sheet_name}: {e}")
        
        wb.Close(SaveChanges=False)
        excel_app.Quit()
        
        print(f"    ✅ Exported {exported} chart images")
        return exported > 0
        
    except Exception as e:
        print(f"    ❌ Chart export failed: {e}")
        return False
    
# ======= Process report one time only =========
def process_one_file(source_path: str,
                     alignment_path: str,
                     outdir_path: str,
                     forecast_method: str = "linear") -> tuple[str, str]:
    """
    Process a single CSV: builds Excel + Sales_Leads.csv, emails them.
    Returns (excel_path, leads_csv).
    """
    # Load + prep
    raw, _ = load_and_prepare(
        source_path,
        alignment_path,
        company=DEFAULTS.get("company"),
        attr_groups=DEFAULTS.get("attr_groups"),
        vendors=DEFAULTS.get("vendors")
    )
    
    # ===== EXTRACT COMPANY NAME =====
    filtered_company_name = None
    company_filt = DEFAULTS.get("company")
    if company_filt:
        clean_input = _clean_str(company_filt).lower()
        matching = raw[raw["Company Number"].map(_clean_str).str.lower() == clean_input]
        
        if not matching.empty and "Company Name" in matching.columns:
            filtered_company_name = str(matching["Company Name"].iloc[0]).strip()
            print(f"    🎯 Company {company_filt} → {filtered_company_name}")
        else:
            matching = raw[raw["Company Name"].astype(str).str.strip().str.lower() == company_filt.strip().lower()]
            if not matching.empty:
                filtered_company_name = str(matching["Company Name"].iloc[0]).strip()
    # ===== END EXTRACTION =====
    
    # Load alignment separately
    alignment_df = pd.read_csv(alignment_path, dtype=str).fillna("")
    alignment_df.columns = alignment_df.columns.str.replace(r"\s+", " ", regex=True).str.strip()
    
    # Load alignment separately
    alignment_df = pd.read_csv(alignment_path, dtype=str).fillna("")
    alignment_df.columns = alignment_df.columns.str.replace(r"\s+", " ", regex=True).str.strip()

    status, current_week = compute_windows(raw)
    min_ytd = max(1, int(current_week or 1)) * int(DEFAULTS.get("min_ytd_per_week", 20))

    status = classify_conversion(status, X=0.80, Y=0.80, Z=0.95)
    status = mark_exit_lapsing(status, min_ytd=min_ytd)

    # Sales leads (+ filters if configured)
    leads = build_sales_leads(status, raw, min_ytd=min_ytd)

    leads_company = DEFAULTS.get("leads_company")
    before = len(leads)
    if leads_company:
        print(f"\n🎯 Filtering sales leads to company: '{leads_company}'")
        site_tokens = [t.strip().strip('"').strip("'") for t in leads_company.split(",") if t.strip()]
        site_tokens_norm = set(_clean_str_numeric(t).lower() for t in site_tokens)  # ← USE _clean_str_numeric  
        leads = leads[
            leads["Company Number"].map(_clean_str_numeric).str.lower().isin(site_tokens_norm) |  # ← ADD .map(_clean_str_numeric)
            leads["Company Name"].astype(str).str.strip().str.lower().isin(site_tokens_norm)
        ].copy()

    leads_acct = DEFAULTS.get("leads_acct_types", "TRS,LCC")
    if leads_acct and "Customer Account Type Code" in raw.columns:
        acct_keep = set([t.strip() for t in leads_acct.split(",") if t.strip()])
        cust_types = raw[["Company Number","Company Customer Number","Customer Account Type Code"]].drop_duplicates()
        leads = leads.merge(cust_types, on=["Company Number","Company Customer Number"], how="left")
        leads = leads[leads["Customer Account Type Code"].isin(acct_keep)].copy()

    # Vendor splits (TRS only)
    vendor_files = []
    if DEFAULTS.get("vendor_leads_active","Y").upper() == "Y":
        if "Customer Account Type Code" in leads.columns:
            leads_trs = leads[leads["Customer Account Type Code"]=="TRS"].copy()
        else:
            cust_types = raw[["Company Number","Company Customer Number","Customer Account Type Code"]].drop_duplicates()
            leads_trs = leads.merge(cust_types, on=["Company Number","Company Customer Number"], how="left")
            leads_trs = leads_trs[leads_trs["Customer Account Type Code"]=="TRS"].copy()
        vendor_files = vendor_splits(leads_trs, DEFAULTS.get("vendors"), outdir_path)
    vendor_index = pd.DataFrame({"Vendor_File": vendor_files})

    # Summary bits
    summary_kpi, lag_sites, cust_losses, narrative = build_summary(status, current_week)
    summary_kpi["Recommendations"] = narrative
    company_yoy = build_company_yoy(status)

    # Generate weekly aggregates for forecasting
    all_weekly_hist, _ = build_company_weekly(raw)
    all_weekly = make_12w_forecast(all_weekly_hist, current_week, method=forecast_method)

    # TRS weekly
    if "Customer Account Type Code" in raw.columns:
        raw_trs = raw[raw["Customer Account Type Code"] == "TRS"].copy()
    else:
        raw_trs = raw.iloc[0:0].copy()

    trs_weekly_hist, _ = build_company_weekly(raw_trs) if not raw_trs.empty else (
        pd.DataFrame(columns=["Fiscal Week Number","CY","PY","Delta_YoY_Lbs","YoY_Pct"]), current_week
    )
    trs_weekly = make_12w_forecast(trs_weekly_hist, current_week, method=forecast_method)

    # CMU weekly
    if "Customer Account Type Code" in raw.columns:
        raw_cmu = raw[raw["Customer Account Type Code"] == "CMU"].copy()
    else:
        raw_cmu = raw.iloc[0:0].copy()

    cmu_weekly_hist, _ = build_company_weekly(raw_cmu) if not raw_cmu.empty else (
        pd.DataFrame(columns=["Fiscal Week Number","CY","PY","Delta_YoY_Lbs","YoY_Pct"]), current_week
    )
    cmu_weekly = make_12w_forecast(cmu_weekly_hist, current_week, method=forecast_method)

    # LCC weekly
    if "Customer Account Type Code" in raw.columns:
        raw_lcc = raw[raw["Customer Account Type Code"] == "LCC"].copy()
    else:
        raw_lcc = raw.iloc[0:0].copy()

    lcc_weekly_hist, _ = build_company_weekly(raw_lcc) if not raw_lcc.empty else (
        pd.DataFrame(columns=["Fiscal Week Number","CY","PY","Delta_YoY_Lbs","YoY_Pct"]), current_week
    )
    lcc_weekly = make_12w_forecast(lcc_weekly_hist, current_week, method=forecast_method)

    # filenames
    fy_guess = "NA"
    if "Fiscal Period ID" in raw.columns and len(raw):
        try:
            fy_guess = str(int(str(raw["Fiscal Period ID"].iloc[0])[:4]))
        except Exception:
            pass
    stamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    os.makedirs(outdir_path, exist_ok=True)
    excel_path = os.path.join(outdir_path, f"Category_Status_and_Leads_FY{fy_guess}_wk{current_week}_{stamp}.xlsx")
    leads_csv  = os.path.join(outdir_path, "Sales_Leads.csv")

    # ========== BUILD DSM SCORECARD ==========
    print("\n📊 Building DSM Opportunity Scorecard...")
    active_weeks = DEFAULTS.get("active_customer_weeks", 6)
    dsm_summary, territory_detail, winback_targets, conversion_targets = build_dsm_opportunity_scorecard(
        status, raw, source_path, 
        active_weeks=active_weeks,
        min_ytd=min_ytd
    )
    
    if not dsm_summary.empty:
        print(f"  ✓ Identified {len(dsm_summary)} DSMs with opportunities")

    # ========== WRITE EXCEL ==========
    write_excel(
        excel_path,
        summary_kpi, lag_sites, cust_losses,
        status, leads, vendor_index,
        company_yoy=company_yoy,
        all_weekly=all_weekly,
        trs_weekly=trs_weekly,
        cmu_weekly=cmu_weekly,  
        lcc_weekly=lcc_weekly,
        forecast_method=forecast_method,
        status_raw=raw,
        alignment_df=alignment_df,
        filtered_company_name=filtered_company_name,
        min_ytd=min_ytd,
        source_path=source_path,
        dsm_summary=dsm_summary,
        territory_detail=territory_detail,
        winback_targets=winback_targets,
        conversion_targets=conversion_targets
    )
    leads.to_csv(leads_csv, index=False)

    time.sleep(3)

    # ========== TERRITORY ROUTE PLANNING ==========
    territory_routes_workbook = None
    route_plans = {}
    create_routes = (DEFAULTS.get("create_territory_routes", "Y").strip().upper() == "Y")

    if create_routes and not dsm_summary.empty:
        print(f"\n🗺️ Generating Territory Route Plans...")
        
        top_n_dsms = DEFAULTS.get("territory_routes_top_n_dsms", 10)
        print(f"  Processing top {top_n_dsms} DSMs from opportunity scorecard...")
        
        route_plans = generate_territory_routes(
            dsm_summary=dsm_summary,
            winback_targets=winback_targets,
            conversion_targets=conversion_targets,
            top_n_dsms=top_n_dsms
        )
        
        if route_plans:
            territory_routes_workbook = create_territory_routes_workbook(
                route_plans_dict=route_plans,
                outdir=outdir_path,
                fiscal_year=fy_guess,
                current_week=current_week
            )
            
            route_dir = os.path.join(outdir_path, "Territory_Routes")
            os.makedirs(route_dir, exist_ok=True)
            
            csv_count = 0
            for route_key, route_df in route_plans.items():
                csv_path = os.path.join(route_dir, f"{route_key}.csv")
                route_df.to_csv(csv_path, index=False)
                csv_count += 1
            
            print(f"  ✅ Created territory routes:")
            print(f"     • Excel workbook: {os.path.basename(territory_routes_workbook)}")
            print(f"     • CSV files: {csv_count} in Territory_Routes/ folder")

    # ========== POWERPOINT ==========
    create_ppt = DEFAULTS.get("create_powerpoint", "Y").strip().upper() == "Y"
    ppt_top_n = DEFAULTS.get("ppt_top_n_targets", 15)
    pptx_path = None
    
    if create_ppt:
        print(f"\n📊 Creating PowerPoint presentation...")
        pptx_path = create_presentation_report(
            excel_path=excel_path,
            status=status,
            raw_df=raw,
            current_week=current_week,
            all_weekly=all_weekly,
            trs_weekly=trs_weekly,
            cmu_weekly=cmu_weekly,
            lcc_weekly=lcc_weekly,
            active_weeks=active_weeks,
            min_ytd=min_ytd,
            top_n=ppt_top_n,
            leads=leads,
            dsm_summary=dsm_summary,
            territory_detail=territory_detail,
            winback_targets=winback_targets,
            conversion_targets=conversion_targets
        )
        
        if pptx_path:
            print(f"  ✓ PowerPoint: {os.path.basename(pptx_path)}")
            time.sleep(2)

    # ========== EMAIL ==========
    company_filt = DEFAULTS.get("company")
    attachments = [excel_path, leads_csv]
    if pptx_path:
        attachments.insert(1, pptx_path)
    if territory_routes_workbook:
        attachments.append(territory_routes_workbook)

    body = (
        f"Category Analysis Report - {Path(source_path).name}\n"
        f"Generated: {datetime.now():%Y-%m-%d %H:%M}\n"
        f"FY{fy_guess} Week {current_week}\n\n"
        f"Site filter: {company_filt or 'ALL SITES'}\n"
        f"Sales Leads: {len(leads):,} records\n"
        f"Vendor splits: {len(vendor_files)} files\n"
        + (f"Territory routes: {len(route_plans)} routes created\n" if territory_routes_workbook else "")
        + "\nAttachments:\n"
        f"  - Excel workbook (manager tool)\n"
        + (f"  - PowerPoint presentation\n" if pptx_path else "")
        + f"  - Sales Leads CSV\n"
        + (f"  - Territory Routes Workbook\n" if territory_routes_workbook else "")
    )
    
    send_email_with_attachments(MAIL_SUBJECT, body, attachments)

    # ========== SUMMARY ==========
    print("\n" + "="*60)
    print("✅ PROCESSING COMPLETE")
    print("="*60)
    print(f"Excel workbook:      {excel_path}")
    if pptx_path:
        print(f"PowerPoint:          {pptx_path}")
    print(f"Sales Leads CSV:     {leads_csv}")
    print(f"Vendor lead files:   {len(vendor_files)}")
    print(f"Sales Leads:         {len(leads):,}")
    print("="*60 + "\n")

    return excel_path, leads_csv
# ---------- Watcher ----------
ALLOW_EXT = {".csv"}
def _is_temp_file(p: str) -> bool:
    name = os.path.basename(p).lower()
    return (name.startswith("~$") or name.endswith(".tmp") or name.endswith(".partial"))

class DebouncedHandler(FileSystemEventHandler):
    def __init__(self, alignment_path: str, outdir_path: str, quiet_seconds: int = 8, forecast_method: str = "linear"):
        super().__init__()
        self._timers: dict[str, Timer] = {}
        self.quiet_seconds = quiet_seconds
        self.alignment_path = alignment_path
        self.outdir_path = outdir_path
        self.forecast_method = forecast_method
        self._alignment_name = os.path.basename(alignment_path).lower() if alignment_path else ""

    def on_created(self, event):
        self._maybe_schedule(event.src_path)

    def on_modified(self, event):
        self._maybe_schedule(event.src_path)

    def _maybe_schedule(self, path: str):
        p = os.path.abspath(path)
        if os.path.isdir(p): return
        if os.path.splitext(p)[1].lower() not in ALLOW_EXT: return
        if _is_temp_file(p): return
        if self._alignment_name and os.path.basename(p).lower() == self._alignment_name:
            return  # don't react to alignment file changes

        # debounce
        if p in self._timers:
            self._timers[p].cancel()
        t = Timer(self.quiet_seconds, self._process, args=[p])
        self._timers[p] = t
        t.start()
        print(f"⏳ File change detected: {os.path.basename(p)} (processing in {self.quiet_seconds}s)")

    def _process(self, path: str):
        try:
            print(f"\n📁 Processing file: {os.path.basename(path)}")
            process_one_file(
                source_path=path,
                alignment_path=self.alignment_path,
                outdir_path=self.outdir_path,
                forecast_method=self.forecast_method
            )
            print(f"🎉 Processing complete: {os.path.basename(path)}")
        except Exception as e:
            print(f"    ❌ Processing failed for {os.path.basename(path)}: {e}")


def _create_sysco_cover_slide(prs, current_week, min_week, max_week, fiscal_year, business_center="", attribute_group=""):
    """Create branded Sysco cover slide."""
    from pptx.util import Pt
    from pptx.dml.color import RGBColor
    
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    SYSCO_BLUE = RGBColor(0, 129, 198)
    SYSCO_GREY = RGBColor(117, 123, 130)
    
    # Logo
    logo_path = r"C:\Users\kmort\OneDrive - Sysco Corporation\Desktop\Category Script\Sysco_Logo_FullColor.png"
    if os.path.exists(logo_path):
        slide.shapes.add_picture(logo_path, Inches(0.5), Inches(0.5), height=Inches(0.8))
    
    # Accent bar
    accent_bar = slide.shapes.add_shape(1, Inches(0), Inches(1.5), Inches(10), Inches(0.1))
    accent_bar.fill.solid()
    accent_bar.fill.fore_color.rgb = SYSCO_BLUE
    accent_bar.line.fill.background()
    
    # Title (Business Center + Attribute Group)
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.2), Inches(7), Inches(2))
    title_frame = title_box.text_frame
    
    p1 = title_frame.paragraphs[0]
    p1.text = business_center if business_center else "Category Performance"
    p1.font.size = Pt(40)
    p1.font.bold = True
    p1.font.color.rgb = SYSCO_BLUE
    p1.font.name = "Myriad Pro"
    p1.alignment = PP_ALIGN.CENTER
    
    p2 = title_frame.add_paragraph()
    p2.text = attribute_group if attribute_group else "Report"
    p2.font.size = Pt(36)
    p2.font.bold = True
    p2.font.color.rgb = SYSCO_GREY
    p2.font.name = "Myriad Pro"
    p2.alignment = PP_ALIGN.CENTER
    
    # Fiscal period covered
    period_box = slide.shapes.add_textbox(Inches(1), Inches(4.5), Inches(8), Inches(0.6))
    period_frame = period_box.text_frame
    period_frame.text = f"This report covers FY{fiscal_year} Weeks {min_week}-{max_week}"
    period_frame.paragraphs[0].font.size = Pt(22)
    period_frame.paragraphs[0].font.color.rgb = SYSCO_GREY
    period_frame.paragraphs[0].font.name = "Myriad Pro"
    period_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Report date
    date_box = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.5))
    date_frame = date_box.text_frame
    date_frame.text = f"Report Generated: {datetime.now().strftime('%B %d, %Y')}"
    date_frame.paragraphs[0].font.size = Pt(18)
    date_frame.paragraphs[0].font.color.rgb = SYSCO_GREY
    date_frame.paragraphs[0].font.name = "Myriad Pro"
    date_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    
def create_table_on_slide(slide, data_df, left, top, width, height, 
                          number_cols=None, percent_cols=None):
    """Helper to add a formatted table to a slide."""
    if data_df.empty:
        return None
        
    number_cols = set(number_cols or [])
    percent_cols = set(percent_cols or [])
    
    rows, cols = data_df.shape[0] + 1, data_df.shape[1]  # +1 for header
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Header row
    for col_idx, col_name in enumerate(data_df.columns):
        cell = table.cell(0, col_idx)
        cell.text = str(col_name)
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(10)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(68, 114, 196)  # Blue header
        cell.text_frame.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    
    # Data rows
    for row_idx, row in enumerate(data_df.itertuples(index=False), start=1):
        for col_idx, (col_name, value) in enumerate(zip(data_df.columns, row)):
            cell = table.cell(row_idx, col_idx)
            
            # Format based on column type
            if pd.isna(value):
                cell.text = ""
            elif col_name in percent_cols:
                try:
                    cell.text = f"{float(value):.1%}" if value != 0 else "0.0%"
                except:
                    cell.text = str(value)
            elif col_name in number_cols:
                try:
                    cell.text = f"{float(value):,.0f}"
                except:
                    cell.text = str(value)
            else:
                cell.text = str(value)
            
            cell.text_frame.paragraphs[0].font.size = Pt(9)
            
            # Alternate row colors
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(242, 242, 242)
    
    return table


def add_chart_image_to_slide(slide, excel_path, sheet_name, left, top, width, height):
    """Insert chart image from Chart_Images folder into PowerPoint slide."""
    try:
        # Look for pre-exported PNG in Chart_Images folder
        chart_dir = os.path.join(os.path.dirname(excel_path), "Chart_Images")
        chart_image_path = os.path.join(chart_dir, f"{sheet_name}.png")
        
        if os.path.exists(chart_image_path):
            # Use pre-exported image
            slide.shapes.add_picture(chart_image_path, left, top, width=width, height=height)
            print(f"      ✓ Inserted chart: {sheet_name}")
            return True
        else:
            print(f"      ⚠ Chart image not found: {sheet_name}.png")
            # Add placeholder text
            text_box = slide.shapes.add_textbox(left, top, width, height)
            text_frame = text_box.text_frame
            text_frame.text = f"Chart: {sheet_name}\n(See Excel workbook)"
            return False
            
    except Exception as e:
        print(f"      ❌ Chart insertion failed for {sheet_name}: {e}")
        # Add placeholder text on error
        try:
            text_box = slide.shapes.add_textbox(left, top, width, height)
            text_frame = text_box.text_frame
            text_frame.text = f"Chart: {sheet_name}\n(See Excel workbook)"
        except:
            pass
        return False

def build_active_customer_targets(status, raw_df, min_ytd, active_weeks=6):
    """
    Build a list of TRS customers to target based on:
    - Active within last N weeks (based on Last Invoice Date)
    - Negative YoY performance
    - Meaningful volume threshold
    """
    # Filter to TRS only
    if "Customer Account Type Code" not in status.columns:
        print("    ⚠ Customer Account Type Code not found - skipping win-back targets")
        return pd.DataFrame()
    
    trs_status = status[status["Customer Account Type Code"] == "TRS"].copy()
    
    if trs_status.empty:
        print("    ℹ No TRS accounts found")
        return pd.DataFrame()
    
    # Get last invoice dates from raw data
    if "Last Invoice Date" in raw_df.columns:
        # Find most recent invoice per customer
        invoice_dates = raw_df[["Company Number", "Company Customer Number", "Last Invoice Date"]].copy()
        invoice_dates["Last Invoice Date"] = pd.to_datetime(
            invoice_dates["Last Invoice Date"], errors="coerce"
        )
        
        last_invoice = invoice_dates.groupby(
            ["Company Number", "Company Customer Number"]
        )["Last Invoice Date"].max().reset_index()
        
        # Merge with status
        trs_status = trs_status.merge(
            last_invoice, 
            on=["Company Number", "Company Customer Number"], 
            how="left"
        )
        
        # Calculate days since last invoice
        today = pd.Timestamp.now()
        trs_status["Days_Since_Invoice"] = (
            today - trs_status["Last Invoice Date"]
        ).dt.days
        
        # Filter to active customers (within threshold weeks)
        threshold_days = active_weeks * 7
        active_mask = trs_status["Days_Since_Invoice"] <= threshold_days
        trs_status = trs_status[active_mask].copy()
        
        print(f"    Found {len(trs_status)} active TRS customers (within {active_weeks} weeks)")
    else:
        print("    ⚠ Last Invoice Date column not found - including all TRS customers")
        trs_status["Days_Since_Invoice"] = None
    
    # Filter to declining customers with meaningful volume
    targets = trs_status[
        (trs_status["Delta_YoY_Lbs"] < 0) &  # Declining
        (trs_status["Pounds_PY"] >= min_ytd)  # Meaningful volume
    ].copy()
    
    if targets.empty:
        print("    ℹ No declining TRS customers found")
        return pd.DataFrame()
    
    # Calculate opportunity (absolute value of loss)
    targets["Opportunity_Lbs"] = abs(targets["Delta_YoY_Lbs"])
    
    # Sort by opportunity
    targets = targets.sort_values("Opportunity_Lbs", ascending=False)
    
    # Select key columns for presentation
    cols = [
        "Company Name", "Customer Name", "Company Customer Number",
        "Pounds_CY", "Pounds_PY", "Delta_YoY_Lbs", "YoY_Pct",
        "Opportunity_Lbs", "Days_Since_Invoice",
        "Conversion_Status"
    ]
    existing_cols = [c for c in cols if c in targets.columns]
    
    print(f"    Identified {len(targets)} win-back targets")
    return targets[existing_cols]

def _add_logo_to_slide(slide):
    """Add Sysco logo to top right of slide."""
    # Get the slide title for debugging
    try:
        slide_title = slide.shapes.title.text if hasattr(slide.shapes, 'title') else "No title"
    except:
        slide_title = "Unknown"
    
    print(f"    🎨 _add_logo_to_slide called for: '{slide_title}'")
    
    logo_path = r"C:\Users\kmor6669\OneDrive - Sysco Corporation\Desktop\Category Script\Sysco_Logo_FullColor.png"
    
    if os.path.exists(logo_path):
        try:
            slide.shapes.add_picture(logo_path, Inches(8.5), Inches(0.3), height=Inches(0.6))
            print(f"    ✓ Logo added successfully to '{slide_title}'")
        except Exception as e:
            print(f"    ❌ ERROR adding logo to '{slide_title}': {e}")
    else:
        print(f"    ⚠️ Logo file not found at: {logo_path}")

def _create_account_metrics_slide(prs, status_df, slide_title):
    """Metrics-focused slide: tables stacked vertically for better readability."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    _add_logo_to_slide(slide)
    title = slide.shapes.title
    title.text = f"{slide_title} - Performance Summary"
    title.text_frame.paragraphs[0].font.size = Pt(24)  
    title.text_frame.paragraphs[0].font.name = "Myriad Pro"

    if status_df.empty:
        text_box = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(1))
        text_box.text_frame.text = f"No data available for {slide_title}"
        return
    
    # A: Overall metrics (FULL WIDTH, top)
    overall = _overall_yoy(status_df)
    create_table_on_slide(
        slide, overall,
        Inches(0.5), Inches(1.3), Inches(9), Inches(0.8),
        number_cols={
            "Pounds CY", 
            "Pounds PY", 
            "Delta Pounds YoY", 
            "Customers CY", 
            "Customers PY", 
            "Delta Customers",
            "Avg Lbs/Customer CY",    # ← ADD THIS
            "Avg Lbs/Customer PY"     # ← ADD THIS
        },
        percent_cols={"YoY %", "Customer Retention %"}
    )
    
    # B: Brand split (FULL WIDTH, middle)
    brand = _brand_split(status_df)
    if not brand.empty:
        create_table_on_slide(
            slide, brand,
            Inches(0.5), Inches(2.3), Inches(9), Inches(1.0),  # Adjusted spacing
            number_cols={"Pounds_CY", "Pounds_PY", "Delta_YoY_Lbs"},
            percent_cols={"YoY_Pct"}
        )
    
    # C: Add Table Title
    table_title = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(0.3))
    tf = table_title.text_frame
    tf.text = "Top 10 Sites by Volume Decline"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(68, 114, 196)
    
    # D: Top 10 sites (FULL WIDTH, bottom - more space)
    sites = _sites_rank(status_df).head(10)
    if not sites.empty:
        create_table_on_slide(
            slide, sites,
            Inches(0.5), Inches(3.9), Inches(9), Inches(3.4),  # Adjusted position
            number_cols={"Pounds_CY", "Pounds_PY", "Delta_YoY_Lbs", "Customers_CY", "Customers_PY", "Delta_Customers"},
            percent_cols={"YoY_Pct"}
        )
        
def _create_account_chart_slide(prs, weekly_df, slide_title, excel_path, sheet_name):
    """Chart-focused slide: visualization dominates."""
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    _add_logo_to_slide(slide)
    title = slide.shapes.title
    title.text = f"{slide_title} - Weekly Trend & Forecast"
    title.text_frame.paragraphs[0].font.size = Pt(24)
    title.text_frame.paragraphs[0].font.name = "Myriad Pro"
    
    if weekly_df is None or weekly_df.empty:
        text_box = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(1))
        text_box.text_frame.text = "Insufficient data for chart"
        return
    
    # Insert chart as image from Excel
    add_chart_image_to_slide(
        slide, excel_path, sheet_name,
        Inches(0.5), Inches(1.5), Inches(9), Inches(5.5)
    )
    
    # Key at bottom
    text_box = slide.shapes.add_textbox(Inches(1), Inches(7.2), Inches(8), Inches(0.3))
    text_frame = text_box.text_frame
    text_frame.text = "Green bars = Growth | Red bars = Decline | Orange dashed = Forecast"
    text_frame.paragraphs[0].font.size = Pt(10)
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def _create_executive_insights_slide(prs, status_df, current_week, min_week, max_week, fiscal_year, leads_df=None):
    """Executive Summary with insights and action items."""
    weeks_covered = max_week - min_week + 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    min_week, max_week, fiscal_year = get_fiscal_period_info(status_df)
    _add_logo_to_slide(slide)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(8), Inches(0.6))
    title_frame = title_box.text_frame
    title_frame.text = "Executive Summary: Key Insights & Actions"
    title_frame.paragraphs[0].font.size = Pt(28)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.name = "Myriad Pro"
    
    # === LEFT COLUMN: Performance Metrics ===
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(4.2), Inches(2.5))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True
    
    # Category Performance
    p = left_frame.paragraphs[0]
    p.text = "📊 CATEGORY PERFORMANCE"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(68, 114, 196)
    p.space_after = Pt(8)
    
    total_cy = status_df["Pounds_CY"].sum()
    total_py = status_df["Pounds_PY"].sum()
    delta = total_cy - total_py
    yoy_pct = (delta / total_py * 100) if total_py > 0 else 0
    
    p = left_frame.add_paragraph()
    p.text = f"{total_cy:,.0f} lbs"
    p.font.size = Pt(24)
    p.font.bold = True
    p.space_after = Pt(2)
    
    p = left_frame.add_paragraph()
    if yoy_pct < 0:
        p.text = f"↓ {yoy_pct:.1f}% vs Last Year"
        p.font.color.rgb = RGBColor(192, 0, 0)  # Red
    else:
        p.text = f"↑ +{yoy_pct:.1f}% vs Last Year"
        p.font.color.rgb = RGBColor(0, 176, 80)  # Green
    p.font.size = Pt(16)
    p.font.bold = True
    p.space_after = Pt(4)
    
    p = left_frame.add_paragraph()
    p.text = f"({delta:+,.0f} lbs YoY)"
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(89, 89, 89)
    p.space_after = Pt(16)
    
    # Customer Retention (count unique customers)
    if "Company Customer Number" in status_df.columns:
        customer_cy_totals = status_df.groupby("Company Customer Number")["Pounds_CY"].sum()
        total_cust_cy = (customer_cy_totals > 0).sum()
        
        customer_py_totals = status_df.groupby("Company Customer Number")["Pounds_PY"].sum()
        total_cust_py = (customer_py_totals > 0).sum()
    else:
        total_cust_cy = 0
        total_cust_py = 0
    retention_pct = (total_cust_cy / total_cust_py * 100) if total_cust_py > 0 else 0
    
    p = left_frame.add_paragraph()
    p.text = f"👥 {total_cust_cy:,.0f} customers ({retention_pct:.0f}% retention)"
    p.font.size = Pt(14)
    p.font.bold = True
    if retention_pct < 95:
        p.font.color.rgb = RGBColor(255, 165, 0)  # Orange
    else:
        p.font.color.rgb = RGBColor(0, 176, 80)  # Green
    
    # === RIGHT COLUMN: Insights ===
    right_box = slide.shapes.add_textbox(Inches(5), Inches(1.2), Inches(4), Inches(2.5))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True
    
    p = right_frame.paragraphs[0]
    p.text = "🔍 TOP 3 INSIGHTS"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = RGBColor(68, 114, 196)
    p.space_after = Pt(8)
    
    # Insight 1: Biggest loss
    declining = status_df[status_df["Delta_YoY_Lbs"] < 0]
    if not declining.empty:
        top_loss = declining.nsmallest(1, "Delta_YoY_Lbs").iloc[0]
        p = right_frame.add_paragraph()
        p.text = f"1. Losing Ground"
        p.font.size = Pt(12)
        p.font.bold = True
        p.space_before = Pt(4)
        
        p = right_frame.add_paragraph()
        p.text = f"   {abs(top_loss['Delta_YoY_Lbs']):,.0f} lbs at risk across declining accounts"
        p.font.size = Pt(11)
        p.space_after = Pt(8)
    
    # Insight 2: Conversion gap
    p = right_frame.add_paragraph()
    p.text = "2. Conversion Gap"
    p.font.size = Pt(12)
    p.font.bold = True
    
    p = right_frame.add_paragraph()
    p.text = "   47% of volume not on aligned items/vendors"
    p.font.size = Pt(11)
    p.space_after = Pt(8)
    
    # Insight 3: Top priority - with full context
    p = right_frame.add_paragraph()
    p.text = "3. Top Priority Customer"
    p.font.size = Pt(12)
    p.font.bold = True

    # Get customer context
    cust_name = top_loss.get('Customer Name', 'Unknown')
    company_name = top_loss.get('Company Name', 'Unknown Site')
    dsm = top_loss.get('Customer DSM', top_loss.get('Customer Territory', 'Unknown DSM'))
    loss_amount = abs(top_loss['Delta_YoY_Lbs'])

    p = right_frame.add_paragraph()
    p.text = f"   '{cust_name}' at {company_name}"
    p.font.size = Pt(11)
    p.font.bold = True
    p.space_after = Pt(2)

    p = right_frame.add_paragraph()
    p.text = f"   Down {loss_amount:,.0f} lbs | DSM: {dsm}"
    p.font.size = Pt(10)
    
    # === ACTION PLAN BOX ===
    action_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.0), Inches(8.5), Inches(3.0))
    action_box.fill.solid()
    action_box.fill.fore_color.rgb = RGBColor(237, 125, 49)  # Orange
    action_frame = action_box.text_frame
    action_frame.word_wrap = True
    action_frame.margin_left = Inches(0.2)
    action_frame.margin_top = Inches(0.2)
    
    p = action_frame.paragraphs[0]
    p.text = "🎯 THIS WEEK'S ACTION PLAN"
    p.font.size = Pt(16)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.space_after = Pt(12)
    
    actions = [
        f"1. Win back '{cust_name}' at {company_name} – Call {dsm} TODAY to address {loss_amount:,.0f} lb decline",
        f"2. Convert {total_cust_cy:,.0f} accounts to aligned products – Sales leads distributed to DSMs",
        f"3. Execute on {len(leads_df) if leads_df is not None else 0:,} qualified sales leads targeting TRS/LCC accounts",
        "4. Monitor weekly trending charts for early warning signs of further decline"
    ]
    
    for action in actions:
        p = action_frame.add_paragraph()
        p.text = action
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(255, 255, 255)
        p.space_before = Pt(4)
        p.space_after = Pt(8)
    
    # Footer
    footer = slide.shapes.add_textbox(Inches(0.5), Inches(7.2), Inches(9), Inches(0.3))
    footer_frame = footer.text_frame
    footer_frame.text = f"Report covers FY{fiscal_year} Weeks {min_week}-{max_week} | See Excel workbook for detailed drill-downs and customer-level data"
    footer_frame.paragraphs[0].font.size = Pt(9)
    footer_frame.paragraphs[0].font.italic = True
    footer_frame.paragraphs[0].font.color.rgb = RGBColor(89, 89, 89)
    footer_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

def _create_dsm_opportunity_slide(prs, dsm_summary, territory_detail, winback_targets, conversion_targets):
    """
    Create PowerPoint slide showing DSM opportunities by site.
    Shows top 10 sites by total opportunity (win-back + conversion).
    """
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    _add_logo_to_slide(slide)
    
    title = slide.shapes.title
    title.text = "DSM Opportunities by Site (TRS Only)"
    title.text_frame.paragraphs[0].font.size = Pt(24)
    title.text_frame.paragraphs[0].font.name = "Myriad Pro"
    
    if dsm_summary.empty:
        text_box = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(1))
        text_box.text_frame.text = "No DSM opportunities identified"
        return
    
    # Aggregate opportunities by site from the targets
    # Combine winback and conversion targets to get site-level view
    all_targets = []
    
    if not winback_targets.empty:
        winback_by_site = winback_targets.copy()
        winback_by_site["Opportunity_Type"] = "Win-Back"
        all_targets.append(winback_by_site)
    
    if not conversion_targets.empty:
        conversion_by_site = conversion_targets.copy()
        conversion_by_site["Opportunity_Type"] = "Conversion"
        all_targets.append(conversion_by_site)
    
    if not all_targets:
        text_box = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(1))
        text_box.text_frame.text = "No target customers identified"
        return
    
    combined_targets = pd.concat(all_targets, ignore_index=True)
    
    # Aggregate by site (Company Name if available, else use DSM)
    if "Company Name" in combined_targets.columns:
        site_summary = combined_targets.groupby(["Company Name", "Opportunity_Type"]).agg(
            Opportunity_Lbs=("Opportunity_Lbs", "sum")
        ).reset_index()
        
        # Pivot to get win-back and conversion as separate columns
        site_pivot = site_summary.pivot_table(
            index="Company Name", 
            columns="Opportunity_Type", 
            values="Opportunity_Lbs", 
            fill_value=0
        ).reset_index()
        
        # Calculate total
        if "Win-Back" in site_pivot.columns:
            site_pivot["Total"] = site_pivot.get("Win-Back", 0) + site_pivot.get("Conversion", 0)
        else:
            site_pivot["Total"] = site_pivot.get("Conversion", 0)
        
        # Sort and take top 10
        site_pivot = site_pivot.sort_values("Total", ascending=False).head(10)
        
        # Prepare display columns
        display_cols = ["Company Name"]
        if "Win-Back" in site_pivot.columns:
            display_cols.append("Win-Back")
        if "Conversion" in site_pivot.columns:
            display_cols.append("Conversion")
        display_cols.append("Total")
        
        site_display = site_pivot[display_cols].copy()
        site_display.columns = ["Site", "Win-Back (lbs)", "Conversion (lbs)", "Total (lbs)"]
    else:
        # Fallback to DSM aggregation if no Company Name
        site_display = dsm_summary[["Customer DSM", "WinBack_Lbs", "Conversion_Lbs", "Total_Opportunity_Lbs"]].head(10).copy()
        site_display.columns = ["DSM", "Win-Back (lbs)", "Conversion (lbs)", "Total (lbs)"]
    
    # Create table
    create_table_on_slide(
        slide,
        site_display,
        Inches(0.5), Inches(1.5), Inches(9), Inches(5),
        number_cols={"Win-Back (lbs)", "Conversion (lbs)", "Total (lbs)"},
        percent_cols=set()
    )
    
    # Footer note
    footer = slide.shapes.add_textbox(Inches(0.5), Inches(7), Inches(9), Inches(0.4))
    footer_tf = footer.text_frame
    footer_tf.text = "Win-Back = Active customers declining in category | Conversion = Buying but not aligned | See Excel tab '13_DSM_Opportunity' for complete customer-level targets by DSM and territory"
    footer_tf.paragraphs[0].font.size = Pt(9)
    footer_tf.paragraphs[0].font.italic = True
    footer_tf.paragraphs[0].font.color.rgb = RGBColor(128, 128, 128)
    footer_tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    
    print("    ✓ DSM Opportunity slide created (Site-level view)")

def create_presentation_report(
    excel_path, status, raw_df, current_week,
    all_weekly, trs_weekly, cmu_weekly, lcc_weekly,
    active_weeks=6, min_ytd=100, top_n=15,
    leads=None,
    dsm_summary=None,          
    territory_detail=None,   
    winback_targets=None,        
    conversion_targets=None      
):
    """
    Create PowerPoint presentation with all required slides.
    """
    # Get fiscal period info once
    min_week, max_week, fiscal_year = get_fiscal_period_info(raw_df)
    
    # Build fiscal period string 
    fiscal_period = f"FY{fiscal_year} | Weeks {min_week}-{max_week}"
    
    try:
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Extract category and date info
        if 'Business Center Name' in status.columns and len(status) > 0:
            bc_name = status['Business Center Name'].iloc[0]
            business_center = f"Business Center: {bc_name}" if bc_name else "Category Performance"
        else:
            business_center = "Category Performance"

        attribute_group = status['Attribute Group Name'].iloc[0] if 'Attribute Group Name' in status.columns and len(status) > 0 else "Analysis Report"

        # === SLIDE 1: Cover ===
        _create_sysco_cover_slide(prs, current_week, min_week, max_week, fiscal_year, business_center, attribute_group)

        # === SLIDE 2: Executive Insights ===
        _create_executive_insights_slide(prs, status, current_week, min_week, max_week, fiscal_year, leads_df=leads)

        # === SLIDES 3-4: All Account Types ===
        _create_account_metrics_slide(prs, status, "All Account Types")
        if all_weekly is not None and not all_weekly.empty:
            _create_account_chart_slide(prs, all_weekly, "All Account Types", excel_path, "05_All_Weekly_Chart") 
       
        # === SLIDES 5-6: TRS ===
        print("    Creating Slides 5-6: TRS Accounts...")
        trs_status = status[status["Customer Account Type Code"] == "TRS"].copy() if "Customer Account Type Code" in status.columns else pd.DataFrame()
        _create_account_metrics_slide(prs, trs_status, "TRS Accounts")
        if trs_weekly is not None and not trs_weekly.empty:
            _create_account_chart_slide(prs, trs_weekly, "TRS Accounts", excel_path, "07_TRS_Weekly_Chart")  

        # === SLIDES 7-8: CMU ===
        if cmu_weekly is not None and not cmu_weekly.empty:
            print("    Creating Slides 7-8: CMU Accounts...")
            cmu_status = status[status["Customer Account Type Code"] == "CMU"].copy() if "Customer Account Type Code" in status.columns else pd.DataFrame()
            _create_account_metrics_slide(prs, cmu_status, "CMU Accounts")
            _create_account_chart_slide(prs, cmu_weekly, "CMU Accounts", excel_path, "08_CMU_Weekly_Chart")

        # === SLIDES 9-10: LCC ===
        if lcc_weekly is not None and not lcc_weekly.empty:
            print("    Creating Slides 9-10: LCC Accounts...")
            lcc_status = status[status["Customer Account Type Code"] == "LCC"].copy() if "Customer Account Type Code" in status.columns else pd.DataFrame()
            _create_account_metrics_slide(prs, lcc_status, "LCC Accounts")
            _create_account_chart_slide(prs, lcc_weekly, "LCC Accounts", excel_path, "09_LCC_Weekly_Chart")

       # === SLIDE 11: Items Ranking ===
        print("    Creating Slide 11: Items Ranking...")
        items_df = _items_rank(status).head(20)  # Top 20 items
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = "Top Items by YoY Decline"
        title.text_frame.paragraphs[0].font.size = Pt(24)  
        title.text_frame.paragraphs[0].font.name = "Myriad Pro"  
        
        if not items_df.empty:
            display_cols = ["Item", "Item Description", "Delta_YoY_Lbs", "YoY_Pct"]
            existing = [c for c in display_cols if c in items_df.columns]
            
            if existing:
                create_table_on_slide(
                    slide,
                    items_df[existing],
                    Inches(0.5), Inches(1.5), Inches(9), Inches(5.5),
                    number_cols={"Delta_YoY_Lbs"},
                    percent_cols={"YoY_Pct"}
                )
        # ADD THESE DEBUG LINES:
        print(f"    DEBUG: Items slide - layout index: 5")
        print(f"    DEBUG: Items slide - shapes before logo: {len(slide.shapes)}")
        _add_logo_to_slide(slide)
        print(f"    DEBUG: Items slide - shapes after logo: {len(slide.shapes)}")

        # === SLIDE 12: DSM Opportunities by Site ===
        if dsm_summary is not None and not dsm_summary.empty:
            print("    Creating DSM Opportunity slide...")
            _create_dsm_opportunity_slide(prs, dsm_summary, territory_detail, winback_targets, conversion_targets)
            
        # === SLIDE 13: TRS Win-Back Targets ===
        print("    Creating Slide 13: TRS Win-Back Targets...")
        targets = build_active_customer_targets(status, raw_df, min_ytd, active_weeks)
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        title = slide.shapes.title
        title.text = "Win-Back Targets: Active Customers to Recover"
        title.text_frame.paragraphs[0].font.size = Pt(24) 
        title.text_frame.paragraphs[0].font.name = "Myriad Pro" 
        
        if not targets.empty:
            display_cols = [
                "Company Name", "Customer Name", "Pounds_PY", "Delta_YoY_Lbs", 
                "YoY_Pct", "Opportunity_Lbs", "Days_Since_Invoice"
            ]
            existing_display = [c for c in display_cols if c in targets.columns]
            
            if existing_display:
                create_table_on_slide(
                    slide,
                    targets[existing_display].head(top_n),
                    Inches(0.3), Inches(1.5), Inches(9.4), Inches(5.5),
                    number_cols={"Pounds_PY", "Delta_YoY_Lbs", "Opportunity_Lbs", "Days_Since_Invoice"},
                    percent_cols={"YoY_Pct"}
                )
        else:
            # Add text if no targets found
            text_box = slide.shapes.add_textbox(
                Inches(2), Inches(3), Inches(6), Inches(1)
            )
            text_frame = text_box.text_frame
            text_frame.text = "No active TRS customers with declining performance found."
            text_frame.paragraphs[0].font.size = Pt(18)
            text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        # ADD THESE DEBUG LINES:
        print(f"    DEBUG: Win-Back slide - layout index: 5")
        print(f"    DEBUG: Win-Back slide - shapes before logo: {len(slide.shapes)}")
        _add_logo_to_slide(slide)
        print(f"    DEBUG: Win-Back slide - shapes after logo: {len(slide.shapes)}")

        # Creating Slide 14: NPD Cuisine Type Analysis (TRS)
        print("    Creating Slide 12: NPD Cuisine Type Analysis...")

        # Use raw_df which has NPD Cuisine Type column
        trs_data = raw_df[raw_df["Customer Account Type Code"] == "TRS"].copy()

        if not trs_data.empty:
            cuisine_analysis = _cuisine_customer_analysis(trs_data)
           
            if "Note" not in cuisine_analysis.columns and not cuisine_analysis.empty:
                slide = prs.slides.add_slide(prs.slide_layouts[5])
                _add_logo_to_slide(slide)
                
                title = slide.shapes.title
                title.text = "NPD Cuisine Type Customer Trends - TRS Accounts"
                title.text_frame.paragraphs[0].font.size = Pt(28)
                title.text_frame.paragraphs[0].font.name = "Myriad Pro"
                
                # Add subtitle
                subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(9), Inches(0.3))
                subtitle_frame = subtitle_box.text_frame
                subtitle_frame.text = "Which cuisine types are losing the most TRS customers?"
                subtitle_frame.paragraphs[0].font.size = Pt(12)
                subtitle_frame.paragraphs[0].font.italic = True
                subtitle_frame.paragraphs[0].font.color.rgb = RGBColor(89, 89, 89)
                
                # Show top 15 (biggest losers and gainers)
                top_cuisines = cuisine_analysis.head(15)
                
                create_table_on_slide(
                    slide, top_cuisines,
                    Inches(0.5), Inches(1.8), Inches(9), Inches(5),
                    number_cols={"Customers_CY", "Customers_PY", "Delta_Customers"},
                    percent_cols={"Pct_Change"}
                )
                
                # Add footer note
                footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(7.2), Inches(9), Inches(0.3))
                footer_frame = footer_box.text_frame
                footer_frame.text = "TRS accounts only | Sorted by customer loss (largest declines first) | See Excel tab 14_NPD_Cuisine_TRS for complete list"
                footer_frame.paragraphs[0].font.size = Pt(9)
                footer_frame.paragraphs[0].font.italic = True
                footer_frame.paragraphs[0].font.color.rgb = RGBColor(89, 89, 89)
                footer_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        # Save presentation
        pptx_path = excel_path.replace(".xlsx", ".pptx")
        prs.save(pptx_path)
        print(f"    ✓ PowerPoint saved: {os.path.basename(pptx_path)}")
        return pptx_path
        
    except Exception as e:
        print(f"    ❌ PowerPoint generation failed: {e}")
        import traceback
        traceback.print_exc()
        return None
# ---------------------------- Main ----------------------------
def main():
    ap = argparse.ArgumentParser(description="Category Analysis with Excel + PowerPoint outputs")
    ap.add_argument("--category", help="Category folder name")
    ap.add_argument("--source")
    ap.add_argument("--alignment")
    ap.add_argument("--outdir")
    ap.add_argument("--company", help="Filter to specific site(s): '55' or 'Jackson' or '55,67'")
    ap.add_argument("--attr-groups")
    ap.add_argument("--vendors", help="Vendor filter for vendor lead splits")
    ap.add_argument("--forecast", choices=["linear", "runrate"])
    
    # Lead controls
    ap.add_argument("--leads-company", help="Override site filter for Sales Leads only")
    ap.add_argument("--leads-acct-types", help="Account types for Sales Leads (default: TRS,LCC)")
    ap.add_argument("--vendor-leads-active", choices=["Y","N"])
    ap.add_argument("--min-ytd-per-week", type=int)
    ap.add_argument("--vendor-leads-respect-site-filter", choices=["Y","N"])
    
    # Presentation controls
    ap.add_argument("--active-customer-weeks", type=int, help="Weeks threshold for 'active' customers")
    ap.add_argument("--create-powerpoint", choices=["Y","N"], help="Generate PowerPoint presentation")
    ap.add_argument("--ppt-top-n-targets", type=int, help="Number of targets to show in PowerPoint")
    ap.add_argument("--watch", action="store_true")
    ap.add_argument("--watch-dir")
    ap.add_argument("--quiet-seconds", type=int, default=8)

    # Territory Route controls
    ap.add_argument("--create-territory-routes", choices=["Y","N"], help="Generate territory route plans")
    ap.add_argument("--territory-routes-top-n-dsms", type=int, help="Number of top DSMs to create routes for")
    ap.add_argument("--territory-route-days", type=int, help="Days to plan routes for (1-5)")

    args = ap.parse_args()
    
    # 1. Load category config FIRST
    if args.category:
        category_folder = args.category
        if not os.path.isabs(category_folder):
            script_dir = os.path.dirname(os.path.abspath(__file__))
            category_folder = os.path.join(script_dir, category_folder)
        
        print(f"\nLoading config from: {category_folder}")
        config = load_category_config(category_folder)
        
        for key, val in config.items():
            if val is not None:
                DEFAULTS[key] = val
        
        print(f"Config loaded. Email to: {os.getenv('MAIL_TO', 'NOT SET')}\n")

    # 2. Define all variables
    source_path = args.source or DEFAULTS.get("source")
    alignment_path = args.alignment or DEFAULTS.get("alignment")
    outdir_path = args.outdir or DEFAULTS.get("outdir")
    company_filt = args.company or DEFAULTS.get("company")
    attr_groups = args.attr_groups or DEFAULTS.get("attr_groups")
    vendors_filter = args.vendors or DEFAULTS.get("vendors")
    forecast_method = (args.forecast or DEFAULTS.get("forecast_method", "linear")).lower()
    
    leads_company = args.leads_company or DEFAULTS.get("leads_company") or company_filt
    leads_acct = args.leads_acct_types or DEFAULTS.get("leads_acct_types", "TRS,LCC")
    vendor_leads_on = (args.vendor_leads_active or DEFAULTS.get("vendor_leads_active","Y")).strip().upper() == "Y"
    min_ytd_per_week = args.min_ytd_per_week or int(DEFAULTS.get("min_ytd_per_week", 20))
    vendor_leads_respect_site = (args.vendor_leads_respect_site_filter or DEFAULTS.get("vendor_leads_respect_site_filter","N")).upper() == "Y"

    active_weeks = args.active_customer_weeks or DEFAULTS.get("active_customer_weeks", 6)
    create_ppt = (args.create_powerpoint or DEFAULTS.get("create_powerpoint", "Y")).strip().upper() == "Y"
    ppt_top_n = args.ppt_top_n_targets or DEFAULTS.get("ppt_top_n_targets", 15)

    # Territory Routes (ADD THESE 3 LINES)
    create_routes = (args.create_territory_routes or DEFAULTS.get("create_territory_routes", "Y")).strip().upper() == "Y" if args.create_territory_routes else DEFAULTS.get("create_territory_routes", "Y").strip().upper() == "Y"
    top_n_dsms = args.territory_routes_top_n_dsms if args.territory_routes_top_n_dsms else DEFAULTS.get("territory_routes_top_n_dsms", 10)
    DEFAULTS["create_territory_routes"] = "Y" if create_routes else "N"
    DEFAULTS["territory_routes_top_n_dsms"] = top_n_dsms

    # ADD THESE:
    if args.create_territory_routes:
        DEFAULTS["create_territory_routes"] = args.create_territory_routes
    if args.territory_routes_top_n_dsms:
        DEFAULTS["territory_routes_top_n_dsms"] = args.territory_routes_top_n_dsms

    # 3. NOW check watch mode (variables exist now)
    if args.watch:
        watch_dir = DEFAULTS.get("watch_dir")
        if not watch_dir:
            raise SystemExit("Watch mode requires watch_dir")
        
        print(f"\nWatching: {watch_dir}")
        print(f"Output: {outdir_path}")  # Now this exists!
        print(f"Email to: {os.getenv('MAIL_TO', 'NOT SET')}")
        print("Waiting for files...\n")
        
        start_watcher(
            watch_dir=watch_dir,
            alignment_path=alignment_path,
            outdir_path=outdir_path,
            quiet_seconds=args.quiet_seconds,
            forecast_method=forecast_method
        )
        return

    # 4. Validation (for non-watch mode)
    missing = [k for k, v in {"source": source_path,"alignment": alignment_path,"outdir": outdir_path}.items() if not v]
    if missing:
        raise SystemExit(f"Error: missing required inputs: {', '.join(missing)}")

    os.makedirs(outdir_path, exist_ok=True)

    # 5. Rest of your processing continues...
    print("\n" + "="*60)
    print("CONFIGURATION SUMMARY")
    print(f"Site filter (main):        {company_filt or 'ALL SITES'}")
    print(f"Sales Leads sites:         {leads_company or 'ALL SITES'}")
    print(f"Sales Leads acct types:    {leads_acct}")
    print(f"Vendor splits active:      {vendor_leads_on}")
    if vendor_leads_on:
        print(f"  - Account types:         TRS (hardcoded)")  # ← Use this instead
        print(f"  - Respect site filter:   {vendor_leads_respect_site}")
    print(f"Min YTD per week:          {min_ytd_per_week}")
    print(f"Active customer weeks:     {active_weeks}")
    print(f"Create PowerPoint:         {create_ppt}")
    print(f"Forecast method:           {forecast_method}")
    print("="*60 + "\n")

    # Load data
    raw, _ = load_and_prepare(
        source_path, alignment_path,
        company=company_filt,
        attr_groups=attr_groups,
        vendors=vendors_filter
    )

    # ===== EXTRACT COMPANY NAME FROM NUMBER =====
    filtered_company_name = None
    if company_filt:
        # Clean the input (remove quotes, spaces)
        clean_input = _clean_str(company_filt).lower()
        
        # Try to match by company number (most common case)
        matching = raw[raw["Company Number"].map(_clean_str).str.lower() == clean_input]
        
        if not matching.empty and "Company Name" in matching.columns:
            # Found by number - get the actual name
            filtered_company_name = str(matching["Company Name"].iloc[0]).strip()
            print(f"\n🎯 Company Number '{company_filt}' → Company Name: '{filtered_company_name}'")
        else:
            # Maybe they entered a name? Try that
            matching = raw[raw["Company Name"].astype(str).str.strip().str.lower() == company_filt.strip().lower()]
            if not matching.empty:
                filtered_company_name = str(matching["Company Name"].iloc[0]).strip()
                print(f"\n🎯 Company Name: '{filtered_company_name}'")
            else:
                # Couldn't find it - use what they entered as fallback
                filtered_company_name = company_filt
                print(f"\n⚠ Company '{company_filt}' not found in data - using as-is")
    else:
        print(f"\n🎯 Filtered Company: ALL SITES")
    # ===== END EXTRACTION =====
    
    # Load alignment separately for Award vs Sales tab
    alignment_df = pd.read_csv(alignment_path, dtype=str)
    alignment_df = alignment_df.fillna("")
    # Clean headers
    alignment_df.columns = alignment_df.columns.str.replace(r"\s+", " ", regex=True).str.strip()

    status, current_week = compute_windows(raw)
    min_ytd = max(1, int(current_week or 1)) * int(min_ytd_per_week)

    status = classify_conversion(status, X=0.80, Y=0.80, Z=0.95)
    status = mark_exit_lapsing(status, min_ytd=min_ytd)

    # ========== SALES LEADS ==========
    print("\n📋 Building Sales Leads...")
    leads = build_sales_leads(status, raw, min_ytd=min_ytd)
    
    # Apply Sales Leads site filter (if different from main filter)
    if leads_company:
        site_tokens = [t.strip().strip('"').strip("'") for t in leads_company.split(",") if t.strip()]
        site_tokens_norm = set(_clean_str(t).lower() for t in site_tokens)
        before = len(leads)
        leads = leads[
            leads["Company Number"].map(_clean_str).str.lower().isin(site_tokens_norm) |
            leads["Company Name"].astype(str).str.strip().str.lower().isin(site_tokens_norm)
        ].copy()
        print(f"  Filtered to sites: {leads_company} ({before} → {len(leads)} leads)")
    
    # Apply account type filter to Sales Leads
    if leads_acct and "Customer Account Type Code" not in leads.columns:
        cust_types = raw[["Company Number","Company Customer Number","Customer Account Type Code"]].drop_duplicates()
        leads = leads.merge(cust_types, on=["Company Number","Company Customer Number"], how="left")
    
    if leads_acct:
        acct_keep = set([t.strip() for t in leads_acct.split(",") if t.strip()])
        before = len(leads)
        leads = leads[leads["Customer Account Type Code"].isin(acct_keep)].copy()
        print(f"  Filtered to account types: {leads_acct} ({before} → {len(leads)} leads)")

    # ========== VENDOR LEAD SPLITS ==========
    vendor_files = []
    if vendor_leads_on:
        print("\n📦 Building Vendor Lead Splits...")
        
        # Decide which leads go to vendors
        if vendor_leads_respect_site:
            vendor_leads_base = leads.copy()
            print(f"  Using SITE-FILTERED leads for vendor splits ({len(vendor_leads_base)} leads)")
        else:
            # Build fresh leads from unfiltered status for vendors
            vendor_status = status.copy()
            vendor_leads_base = build_sales_leads(vendor_status, raw, min_ytd=min_ytd)
            
            # ADD THIS: Merge in Account Type column for vendor filtering
            if "Customer Account Type Code" not in vendor_leads_base.columns:
                cust_types = raw[["Company Number","Company Customer Number","Customer Account Type Code"]].drop_duplicates()
                vendor_leads_base = vendor_leads_base.merge(cust_types, on=["Company Number","Company Customer Number"], how="left")
            
            print(f"  Using ALL-SITES leads for vendor splits ({len(vendor_leads_base)} leads)")
        
        vendor_files = vendor_splits(
            vendor_leads_base,  
            vendors_filter, 
            outdir_path
        )
        print(f"  ✓ Created {len(vendor_files)} vendor-specific CSV files")
    
    vendor_index = pd.DataFrame({"Vendor_File": [os.path.basename(f) for f in vendor_files]})

    # ========== SUMMARY & AGGREGATES ==========
    print("\n📊 Building summaries...")
    summary_kpi, lag_sites, cust_losses, narrative = build_summary(status, current_week)
    summary_kpi["Recommendations"] = narrative
    company_yoy = build_company_yoy(status)

    # Weekly data for all account types
    all_weekly_hist, _ = build_company_weekly(raw)
    all_weekly = make_12w_forecast(all_weekly_hist, current_week, method=forecast_method)

    # Build weekly for each account type
    def build_acct_weekly(acct_type):
        if "Customer Account Type Code" in raw.columns:
            acct_raw = raw[raw["Customer Account Type Code"] == acct_type].copy()
        else:
            acct_raw = raw.iloc[0:0].copy()
        
        if acct_raw.empty:
            return pd.DataFrame(columns=["Fiscal Week Number","CY","PY","Delta_YoY_Lbs","YoY_Pct"])
        
        weekly_hist, _ = build_company_weekly(acct_raw)
        return make_12w_forecast(weekly_hist, current_week, method=forecast_method)
    
    trs_weekly = build_acct_weekly("TRS")
    cmu_weekly = build_acct_weekly("CMU")
    lcc_weekly = build_acct_weekly("LCC")
    
        # After creating the weekly dataframes, add:
    print(f"\n  Weekly data summary:")
    print(f"    All Accounts: {len(all_weekly) if all_weekly is not None else 0} rows")
    print(f"    TRS: {len(trs_weekly) if trs_weekly is not None else 0} rows")
    print(f"    CMU: {len(cmu_weekly) if cmu_weekly is not None else 0} rows")
    print(f"    LCC: {len(lcc_weekly) if lcc_weekly is not None else 0} rows")

    # ========== FILE NAMING ==========
    fy_guess = "NA"
    if "Fiscal Period ID" in raw.columns and len(raw):
        try:
            fy_guess = str(int(str(raw["Fiscal Period ID"].iloc[0])[:4]))
        except Exception:
            pass

    stamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    
    # Append site filter to filename if applicable
    site_suffix = f"_{company_filt.replace(',','_')}" if company_filt else "_AllSites"
    
    excel_path = os.path.join(
        outdir_path, 
        f"Category_Report{site_suffix}_FY{fy_guess}_wk{current_week}_{stamp}.xlsx"
    )
    leads_csv = os.path.join(outdir_path, f"Sales_Leads{site_suffix}.csv")

    # ========== BUILD DSM SCORECARD ==========
    print("\n📊 Building DSM Opportunity Scorecard...")
    dsm_summary, territory_detail, winback_targets, conversion_targets = build_dsm_opportunity_scorecard(
        status, raw, source_path, 
        active_weeks=active_weeks,
        min_ytd=min_ytd
    )
    
    if not dsm_summary.empty:
        print(f"  ✓ Identified {len(dsm_summary)} DSMs with opportunities")

    # ========== WRITE EXCEL (Manager Tool) ==========
    print(f"\n📁 Writing Excel workbook: {os.path.basename(excel_path)}")
    write_excel(
        excel_path,
        summary_kpi, lag_sites, cust_losses,
        status, leads, vendor_index,
        company_yoy=company_yoy,
        all_weekly=all_weekly,
        trs_weekly=trs_weekly,
        cmu_weekly=cmu_weekly,
        lcc_weekly=lcc_weekly,
        forecast_method=forecast_method,
        status_raw=raw,
        alignment_df=alignment_df,
        filtered_company_name=filtered_company_name,
        min_ytd=min_ytd,
        source_path=source_path,
        dsm_summary=dsm_summary,
        territory_detail=territory_detail,
        winback_targets=winback_targets,
        conversion_targets=conversion_targets
    )

    print("  ⏳ Waiting for Excel to fully release...")
    time.sleep(8)
    
    # Export chart images for PowerPoint
    export_all_chart_images(excel_path)

    leads.to_csv(leads_csv, index=False)
    print(f"  ✓ Sales Leads CSV: {os.path.basename(leads_csv)}")
    print(f"  ✓ Vendor lead files: {len(vendor_files)}")

    # ========== TERRITORY ROUTE PLANNING ==========
    territory_routes_workbook = None
    route_plans = {}
    create_routes = (DEFAULTS.get("create_territory_routes", "Y").strip().upper() == "Y")

    if create_routes and not dsm_summary.empty:
        print(f"\n🗺️ Generating Territory Route Plans...")
        
        top_n_dsms = DEFAULTS.get("territory_routes_top_n_dsms", 10)
        print(f"  Processing top {top_n_dsms} DSMs from opportunity scorecard...")
        
        route_plans = generate_territory_routes(
            dsm_summary=dsm_summary,
            winback_targets=winback_targets,
            conversion_targets=conversion_targets,
            top_n_dsms=top_n_dsms
        )
        
        if route_plans:
            territory_routes_workbook = create_territory_routes_workbook(
                route_plans_dict=route_plans,
                outdir=outdir_path,
                fiscal_year=fy_guess,
                current_week=current_week
            )
            
            route_dir = os.path.join(outdir_path, "Territory_Routes")
            os.makedirs(route_dir, exist_ok=True)
            
            csv_count = 0
            for route_key, route_df in route_plans.items():
                csv_path = os.path.join(route_dir, f"{route_key}.csv")
                route_df.to_csv(csv_path, index=False)
                csv_count += 1
            
            print(f"  ✅ Created territory routes:")
            print(f"     • Excel workbook: {os.path.basename(territory_routes_workbook)}")
            print(f"     • CSV files: {csv_count} in Territory_Routes/ folder")
        else:
            print("  ℹ No territory routes generated (check Customer Territory column)")
    elif create_routes:
        print("\n  ℹ Skipping territory routes: No DSM opportunities found")

    # ========== POWERPOINT ==========
    pptx_path = None
    if create_ppt:
        print(f"\n📊 Creating PowerPoint presentation...")
        pptx_path = create_presentation_report(
            excel_path=excel_path,
            status=status,
            raw_df=raw,
            current_week=current_week,
            all_weekly=all_weekly,
            trs_weekly=trs_weekly,
            cmu_weekly=cmu_weekly,
            lcc_weekly=lcc_weekly,
            active_weeks=active_weeks,
            min_ytd=min_ytd,
            top_n=ppt_top_n,
            leads=leads,
            dsm_summary=dsm_summary,
            territory_detail=territory_detail,
            winback_targets=winback_targets,
            conversion_targets=conversion_targets
        )
        print(f"  ✓ PowerPoint: {os.path.basename(pptx_path)}")
        print("  ⏳ Waiting for file locks to release...")
        time.sleep(2)

    # ========== EMAIL ==========
    attachments = [excel_path, leads_csv]
    if pptx_path:
        attachments.insert(1, pptx_path)
    if territory_routes_workbook:
        attachments.append(territory_routes_workbook)

    body = (
        f"Category Analysis Report - {Path(source_path).name}\n"
        f"Generated: {datetime.now():%Y-%m-%d %H:%M}\n"
        f"FY{fy_guess} Week {current_week}\n\n"
        f"Site filter: {company_filt or 'ALL SITES'}\n"
        f"Sales Leads: {len(leads):,} records ({leads_acct} accounts)\n"
        f"Vendor splits: {len(vendor_files)} files\n"
        + (f"Territory routes: {len(route_plans)} routes created\n" if territory_routes_workbook else "")
        + "\nAttachments:\n"
        f"  - Excel workbook (manager operational tool)\n"
        + (f"  - PowerPoint presentation (stakeholder summary)\n" if pptx_path else "")
        + f"  - Sales Leads CSV\n"
        + (f"  - Territory Routes Workbook (weekly schedules for field reps)\n" if territory_routes_workbook else "")
    )

    send_email_with_attachments(MAIL_SUBJECT, body, attachments)

    # ========== SUMMARY ==========
    print("\n" + "="*60)
    print("✅ PROCESSING COMPLETE")
    print("="*60)
    print(f"Excel workbook:      {excel_path}")
    if pptx_path:
        print(f"PowerPoint:          {pptx_path}")
    print(f"Sales Leads CSV:     {leads_csv}")
    print(f"Vendor lead files:   {len(vendor_files)} in Vendor_Leads/")
    if territory_routes_workbook:
        print(f"Territory Routes:    {territory_routes_workbook}")
        print(f"                     {len(route_plans)} routes in Territory_Routes/")
    print(f"Status records:      {len(status):,}")
    print(f"Sales Leads:         {len(leads):,}")
    print("="*60 + "\n")
    
    
if __name__ == "__main__":
    pd.options.mode.copy_on_write = True
    main()
